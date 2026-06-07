import os
import sys
import logging
import types
import warnings
import subprocess
import shutil
import cv2
import uuid
import numpy as np
import fitz  # библиотека PyMuPDF
from pptx import Presentation
import concurrent.futures
from concurrent.futures import ThreadPoolExecutor, as_completed
import torch
from llama_index.core.schema import TextNode
import base64
import requests
import json
import time
import gc
import threading
from llama_index.core.node_parser import SentenceSplitter
from src.gguf_direct import detect_model_family, get_gguf_llm, unload_all_models
import config

logger = logging.getLogger(__name__)

# F-fix #11: модульный кеш WhisperX-модели.
# Ключ — (model_name, device, compute_type). Без кеша load_model занимает
# ~30 сек и загружает ~1.5 GB в VRAM при КАЖДОМ аудио/видео файле.
# В batch'е из 5 mp3 это лишние 2 минуты загрузок. С кешем — 1 раз за batch.
# Сбрасывается через unload_whisper_model() (вызывается из main.py в конце batch'а).
_whisper_model_cache: dict = {}
_whisper_lock = threading.Lock()


# F-fix #15: общий requests.Session с HTTPAdapter.
# В OCR batch'е делается 60+ vision-запросов подряд. Без Session каждый
# запрос открывает новый TCP-коннект (50-100 мс handshake). С Session —
# переиспользуется. Экономия: 1-3 секунды на 60-кадровое видео.
_http_session = requests.Session()
_http_session.mount("http://", requests.adapters.HTTPAdapter(pool_connections=config.HTTP_POOL_SIZE_INGEST, pool_maxsize=config.HTTP_POOL_SIZE_INGEST))
_http_session.mount("https://", requests.adapters.HTTPAdapter(pool_connections=config.HTTP_POOL_SIZE_INGEST, pool_maxsize=config.HTTP_POOL_SIZE_INGEST))


def get_or_load_whisper(model_name: str = "medium", device: str = "cuda", compute_type: str = "int8"):
    """Возвращает кешированную WhisperX-модель или грузит и кладёт в кеш.

    Потокобезопасно (lock). При первом вызове — load_model (~30 сек, 1.5 GB VRAM).
    При последующих — возврат ссылки на тот же объект.
    """
    import whisperx
    key = (model_name, device, compute_type)
    with _whisper_lock:
        if key in _whisper_model_cache:
            return _whisper_model_cache[key]
        _safe_print(f"[WhisperX] Загрузка модели {model_name} ({device}, {compute_type})...")
        model = whisperx.load_model(model_name, device, compute_type=compute_type)
        _whisper_model_cache[key] = model
        return model


def unload_whisper_model():
    """Выгружает ВСЕ кешированные WhisperX-модели и освобождает VRAM.

    Вызывается из main.py после последнего аудио/видео файла в batch'е.
    """
    with _whisper_lock:
        if not _whisper_model_cache:
            return
        _safe_print(f"[WhisperX] Выгрузка {len(_whisper_model_cache)} кешированных моделей...")
        _whisper_model_cache.clear()
        gc.collect()
        if torch.cuda.is_available():
            torch.cuda.empty_cache()

# Подавляем шумные предупреждения
warnings.filterwarnings("ignore", message="Module 'speechbrain")
warnings.filterwarnings("ignore", message="torchcodec is not installed")
warnings.filterwarnings("ignore", message="TensorFloat-32")
warnings.filterwarnings("ignore", message=".*speechbrain.*deprecated", category=UserWarning)
warnings.filterwarnings("ignore", message=".*Lightning automatically upgraded.*")

import logging
logging.getLogger("lightning.pytorch.utilities.migration").setLevel(logging.ERROR)
logging.getLogger("lightning.pytorch").setLevel(logging.ERROR)
logging.getLogger("whisperx").setLevel(logging.WARNING)

# --- Фикс: inspect.stack() ---
import inspect as _inspect_module
_orig_getmodule = _inspect_module.getmodule
def _safe_getmodule(obj, filename=None):
    try: return _orig_getmodule(obj, filename)
    except Exception: return None
_inspect_module.getmodule = _safe_getmodule

class IngestionCancelled(Exception):
    """Поднимается, когда пользователь запросил отмену обработки файла."""
    pass

# Глобальный реестр активных subprocess-ов по notebook_id.
# Нужен, чтобы cancel мог их мгновенно убить (ffmpeg, whisperx и т.д.),
# иначе блокирующий вызов может идти минуты после нажатия кнопки.
_active_subprocesses: dict = {}

def register_subprocess(notebook_id, popen):
    _active_subprocesses.setdefault(notebook_id, []).append(popen)

def unregister_subprocess(notebook_id, popen):
    lst = _active_subprocesses.get(notebook_id)
    if lst and popen in lst:
        try: lst.remove(popen)
        except Exception: pass
    if lst is not None and not lst:
        _active_subprocesses.pop(notebook_id, None)

def kill_subprocesses(notebook_id):
    """Немедленно убивает все subprocess-ы, зарегистрированные для этого блокнота."""
    procs = _active_subprocesses.pop(notebook_id, [])
    for p in procs:
        try:
            if p.poll() is None:
                p.terminate()
        except Exception:
            pass
    return len(procs)

# --- Фикс для Windows DLL ---
try:
    import torch
    lib_dir = os.path.join(os.path.dirname(torch.__file__), 'lib')
    if os.path.exists(lib_dir): os.add_dll_directory(lib_dir)
except Exception:
    pass

import whisperx

def cleanup_gpu():
    """Принудительная очистка всей видеопамяти перед тяжелыми задачами."""
    try:
        from src.rag_pipeline import unload_rag_models
        unload_rag_models(hard=False)
        # Мы НЕ выгружаем GGUF модели здесь, так как get_gguf_llm сам решит, 
        # нужно ли перезапускать сервер или использовать текущий.
        import gc, torch
        gc.collect()
        if torch.cuda.is_available():
            torch.cuda.empty_cache()
        print("[GPU] Память полностью очищена для анализа.")
    except Exception as e:
        print(f"[GPU] Ошибка при очистке: {e}")

def format_seconds(s):
    h = int(s // 3600); m = int((s % 3600) // 60); sec = int(s % 60)
    return f"{h}:{m:02d}:{sec:02d}" if h > 0 else f"{m}:{sec:02d}"

def get_image_base64(image_path):
    with open(image_path, "rb") as image_file: return base64.b64encode(image_file.read()).decode("utf-8")

def get_vision_url(llm_settings, progress_cb=None):
    """Ленивая инициализация Vision-сервера только когда он реально нужен."""
    if not llm_settings or not llm_settings.get("use_gguf_direct"):
        return None
        
    v_model = llm_settings.get("vision_model_path") or llm_settings.get("gguf_model_path")
    if not v_model:
        return None
        
    try:
        if progress_cb: progress_cb(60, "Запуск Vision-сервера (ленивая загрузка)...")
        from src.ingestion import cleanup_gpu
        cleanup_gpu()
        
        v_mmproj = llm_settings.get("vision_mmproj_path") or llm_settings.get("gguf_mmproj_path")
        g_path = config.resolve_model_path(v_model)
        m_path = config.resolve_model_path(v_mmproj)
        v_ctx = int(llm_settings.get("vision_ctx_size") or config.GGUF_CTX_SIZE)
        v_gl = int(llm_settings.get("vision_gpu_layers") or -1)
        v_b = int(llm_settings.get("vision_batch_size") or 512)
        v_ub = int(llm_settings.get("vision_ubatch_size") or 256)
        v_fa = llm_settings.get("vision_flash_attn") == "true"
        v_kv = int(llm_settings.get("vision_kv_quant") or 2)
        v_conc = int(llm_settings.get("vision_concurrency") or config.VISION_CONCURRENCY)
        v_mtp = bool(llm_settings.get("vision_mtp_enabled", False))

        return get_gguf_llm(
            gguf_path=g_path, mmproj_path=m_path,
            ctx_size=v_ctx, gpu_layers=v_gl, n_batch=v_b, n_ubatch=v_ub, flash_attn=v_fa,
            type_k=v_kv, type_v=v_kv,
            n_parallel=v_conc,
            mtp_enabled=v_mtp,
            custom_args=["--reasoning", "off", "--reasoning-format", "none", "--reasoning-budget", "0", "--no-context-shift"]
        )
    except Exception as e:
        print(f"[Vision] Ошибка ленивого запуска: {e}")
        return None

def describe_image_with_lmstudio(image_path, llm_settings=None, existing_llm_url=None):
    def _clean_think_tags(text):
        import re
        # Удаляем блоки рассуждений и любые системные токены типа <|...|>
        text = re.sub(r'<think>.*?</think>', '', text, flags=re.DOTALL)
        text = re.sub(r'<\|.*?\|>', '', text)
        text = re.sub(r'<start_of_turn>|<end_of_turn>', '', text)
        return text.strip()

    prompt = """Проведи детальный анализ изображения. 

1. ГРАФИКА: 
   - Если есть схемы, диаграммы или таблицы, опиши их структуру максимально подробно. 
   - Используй таблицу (Узлы, Связи, Пояснения) для сложных схем.
   - Описывай каждый элемент и связь только ОДИН раз. Не дублируй информацию.
   - Укажи назначение ключевых обозначений, если они присутствуют.

2. ТЕКСТ: 
   - Выполни OCR только того текста, который является текстовым блоком (заголовки, подписи к рисункам, абзацы, списки, вопросы). 
   - НЕ описывай элементы графики в этом разделе, если они уже были детально описаны в п.1.
   - Сохраняй исходную структуру текста.

3. КОНТЕКСТ: 
   - В 1-2 предложениях опиши общую суть и назначение страницы.

Пиши технически точно, лаконично, без лишних вводных фраз и пояснений процесса."""

    # Если передан URL запущенного сервера llama-server
    if existing_llm_url:
        for attempt in range(2):
            try:
                v_temp = float(llm_settings.get("vision_temperature") or config.VISION_TEMPERATURE)
                v_max = int(llm_settings.get("vision_max_tokens") or 4096)
                v_r_pen = float(llm_settings.get("vision_repeat_penalty") or config.VISION_REPEAT_PENALTY)
                v_top_p = float(llm_settings.get("vision_top_p") or config.VISION_TOP_P)
                v_min_p = float(llm_settings.get("vision_min_p") or config.VISION_MIN_P)
                v_pres = float(llm_settings.get("vision_presence_penalty") or 0.0)
                v_freq = float(llm_settings.get("vision_frequency_penalty") or 0.0)
                
                payload = {
                    "messages": [{"role":"user","content":[{"type":"text","text":prompt},{"type":"image_url","image_url":{"url":f"data:image/jpeg;base64,{get_image_base64(image_path)}"}}]}],
                    "temperature": v_temp, 
                    "max_tokens": v_max,
                    "repeat_penalty": v_r_pen,
                    "top_p": v_top_p,
                    "min_p": v_min_p,
                    "presence_penalty": v_pres,
                    "frequency_penalty": v_freq
                }
                r = _http_session.post(f"{existing_llm_url}/v1/chat/completions", json=payload, timeout=300)
                if r.status_code == 200:
                    res = r.json()
                    if "choices" in res:
                        ans = res["choices"][0]["message"]["content"]
                        reason = res["choices"][0].get("finish_reason")
                        ans = _clean_think_tags(ans)
                        print(f"[Ingestion] Описание получено ({len(ans)} симв.). Причина завершения: {reason}")
                        return ans
                elif r.status_code == 500:
                    print(f"[Ingestion] GGUF 500 (попытка {attempt+1}). Повтор...")
                    time.sleep(2)
                    continue
                else:
                    print(f"[Ingestion] Ошибка GGUF {r.status_code}: {r.text}")
            except Exception as e:
                print(f"[Ingestion] Ошибка запроса (попытка {attempt+1}): {e}")
                time.sleep(1)
        return "Ошибка анализа после всех попыток"

    # Резервный вариант через LM Studio или другой OpenAI API
    api_url = (llm_settings.get("llm_url") if llm_settings else None) or config.LM_STUDIO_URL
    api_key = (llm_settings.get("llm_api_key") if llm_settings else None) or config.LLM_DEFAULT_API_KEY
    model_name = (llm_settings.get("llm_model") if llm_settings else None) or config.LLM_DEFAULT_MODEL
    try:
        v_temp = float(llm_settings.get("vision_temperature") or 0.2)
        payload = {
            "model": model_name, 
            "messages": [{"role":"user","content":[{"type":"text","text":prompt},{"type":"image_url","image_url":{"url":f"data:image/jpeg;base64,{get_image_base64(image_path)}"}}]}], 
            "temperature": v_temp
        }
        r = _http_session.post(f"{api_url.rstrip('/')}/chat/completions", headers={"Authorization":f"Bearer {api_key}"}, json=payload, timeout=30)
        ans = r.json()["choices"][0]["message"]["content"]
        return _clean_think_tags(ans)
    except Exception as e:
        logger.warning(f"Ошибка резервного Vision через LM Studio: {e}")
        return "Изображение без описания."

def save_high_res_frame(video_path, time_sec, output_path):
    try:
        from imageio_ffmpeg import get_ffmpeg_exe
        cmd = [get_ffmpeg_exe(), "-y", "-hwaccel", "cuda", "-ss", str(time_sec), "-i", video_path, "-vframes", "1", "-vf", "scale=-2:720", "-q:v", "4", output_path]
        subprocess.run(cmd, capture_output=True)
    except Exception as e:
        logger.warning(f"Ошибка FFmpeg при сохранении кадра: {e}")

def process_audio_video(file_path, images_dir, is_video=False, progress_cb=None, llm_settings=None, cancel_check=None, notebook_id=None, keep_vision_alive=False, keep_whisper_alive=False):
    def _is_cancelled():
        return bool(cancel_check and cancel_check())

    file_name = os.path.basename(file_path)
    _safe_print(f"[process_audio_video] Начало: {file_name}, is_video={is_video}")
    def prog(pct, msg):
        try: print(f"  [{pct}%] {msg}")
        except Exception: pass
        if progress_cb: progress_cb(pct, msg)

    nodes = []
    transcript_data = []
    frame_data = []

    # 1. ТРАНСКРИБАЦИЯ
    # F-fix #11: используем модульный кеш. В batch'е WhisperX грузится 1 раз,
    # а не N раз. keep_whisper_alive=False → выгрузить после текущего файла.
    prog(15, "Транскрибация речи (WhisperX)...")
    device = "cuda" if torch.cuda.is_available() else "cpu"
    try:
        import whisperx
        model = get_or_load_whisper("medium", device, "int8")
        audio = whisperx.load_audio(file_path)
        result = model.transcribe(audio, batch_size=16)
        for seg in result.get('segments', []):
            transcript_data.append({"start": seg['start'], "end": seg['end'], "text": seg['text'].strip()})
    except Exception as e: print(f"WhisperX error: {e}")
    finally:
        if not keep_whisper_alive:
            unload_whisper_model()

    if transcript_data:
        chunk_text = ""; chunk_start = 0
        for i, seg in enumerate(transcript_data):
            if not chunk_text: chunk_start = seg["start"]
            chunk_text += f"[{seg['start']:.1f}s] {seg['text']} "
            if (seg["end"] - chunk_start > 60) or (i == len(transcript_data) - 1):
                nodes.append(TextNode(text=f"Транскрипт {file_name}:\n{chunk_text.strip()}", metadata={"file_name":file_name, "start":chunk_start}))
                chunk_text = ""
    prog(60, "Транскрибация завершена")

    # 2. АНАЛИЗ ВИДЕО
    if is_video:
        prog(62, "Анализ изменений в видео...")
        cap = cv2.VideoCapture(file_path)
        try:
            fps = cap.get(cv2.CAP_PROP_FPS) or 25
            total_frames = int(cap.get(cv2.CAP_PROP_FRAME_COUNT))
            duration_sec = total_frames / fps if fps > 0 else 0
        finally: cap.release()

        from imageio_ffmpeg import get_ffmpeg_exe
        ffmpeg = get_ffmpeg_exe()
        
        PIXEL_THR = 15; UPDATE_PCT = 0.002; NEW_SLIDE_PCT = 0.04; MOTION_PCT = 0.002
        STABLE_WAIT_SEC = 3.0; CHECK_STEP_SEC = 1.0; COMPARE_SIZE = (320, 180)
        
        cmd = [ffmpeg, "-hide_banner", "-loglevel", "error", "-hwaccel", "cuda", "-hwaccel_output_format", "cuda", "-i", file_path, 
               "-vf", f"fps=1/{CHECK_STEP_SEC},scale_cuda={COMPARE_SIZE[0]}:{COMPARE_SIZE[1]}:format=yuv420p,hwdownload,format=yuv420p,format=bgr24", 
               "-f", "image2pipe", "-vcodec", "rawvideo", "-pix_fmt", "bgr24", "pipe:1"]
        
        process = subprocess.Popen(cmd, stdout=subprocess.PIPE, bufsize=10**8)
        if notebook_id is not None:
            register_subprocess(notebook_id, process)
        frame_list = []
        try:
            prev_saved_thumb = None; last_seen_thumb = None; stable_since_sec = 0; current_sec = 0
            chunk_size = COMPARE_SIZE[0] * COMPARE_SIZE[1] * 3
            cancel_check_every = 30  # проверяем cancel каждые ~30 кадров (~30с видео)
            frame_counter = 0
            while True:
                if frame_counter % cancel_check_every == 0 and _is_cancelled():
                    print(f"[Ingestion] Отмена во время извлечения кадров видео ({format_seconds(current_sec)})")
                    try: process.terminate()
                    except Exception: pass
                    raise IngestionCancelled(f"Cancelled during frame extraction at {format_seconds(current_sec)}")
                frame_counter += 1
                raw_frame = process.stdout.read(chunk_size)
                if not raw_frame or len(raw_frame) != chunk_size: break
                thumb = np.frombuffer(raw_frame, dtype='uint8').reshape((COMPARE_SIZE[1], COMPARE_SIZE[0], 3))
                if last_seen_thumb is None:
                    last_seen_thumb = thumb; stable_since_sec = current_sec
                    img_path = os.path.join(images_dir, f"v_{uuid.uuid4().hex[:6]}.jpg")
                    save_high_res_frame(file_path, current_sec, img_path)
                    frame_list.append((img_path, current_sec)); prev_saved_thumb = thumb
                else:
                    diff_motion = cv2.absdiff(thumb, last_seen_thumb)
                    motion_pct = float(np.sum(diff_motion > PIXEL_THR)) / diff_motion.size
                    if motion_pct >= MOTION_PCT: stable_since_sec = current_sec
                    else:
                        if current_sec - stable_since_sec >= STABLE_WAIT_SEC:
                            if prev_saved_thumb is not None:
                                diff_saved = cv2.absdiff(thumb, prev_saved_thumb)
                                saved_pct = float(np.sum(diff_saved > PIXEL_THR)) / diff_saved.size
                                if saved_pct >= UPDATE_PCT:
                                    img_path = os.path.join(images_dir, f"v_{uuid.uuid4().hex[:6]}.jpg")
                                    save_high_res_frame(file_path, current_sec, img_path)
                                    if saved_pct >= NEW_SLIDE_PCT: frame_list.append((img_path, current_sec))
                                    else:
                                        if frame_list: 
                                            try: os.remove(frame_list[-1][0])
                                            except Exception: pass
                                            frame_list[-1] = (img_path, current_sec)
                                    prev_saved_thumb = thumb
                            stable_since_sec = current_sec
                last_seen_thumb = thumb; current_sec += CHECK_STEP_SEC
                if int(current_sec) % 5 == 0:
                    prog(62 + int((current_sec / duration_sec) * 3) if duration_sec > 0 else 62, f"Анализ видео: {format_seconds(current_sec)} / {format_seconds(duration_sec)}")
        finally:
            process.stdout.close(); process.terminate(); process.wait()
            if notebook_id is not None:
                unregister_subprocess(notebook_id, process)

        # 3. ОПИСАНИЕ КАДРОВ
        n = len(frame_list)
        shared_llm_url = None
        if n > 0:
            prog(65, f"Запуск Vision-сервера для описания {n} кадров...")
            shared_llm_url = get_vision_url(llm_settings, progress_cb=prog)

        prog(65, f"Описание {n} кадров ({'параллельно' if (int(llm_settings.get('vision_concurrency') or config.VISION_CONCURRENCY)) > 1 else 'последовательно'})...")
        v_conc = int(llm_settings.get("vision_concurrency") or config.VISION_CONCURRENCY)
        
        from concurrent.futures import ThreadPoolExecutor
        with ThreadPoolExecutor(max_workers=v_conc) as executor:
            # Создаем список задач
            futures = []
            for idx, (path, t) in enumerate(frame_list):
                futures.append(executor.submit(describe_image_with_lmstudio, path, llm_settings, shared_llm_url))
            
            # Собираем результаты по мере готовности для обновления прогресса
            # Используем большой chunk_size, чтобы описания не разрывались
            splitter = SentenceSplitter(chunk_size=2048, chunk_overlap=128)
            try:
                for idx, future in enumerate(futures):
                    if _is_cancelled():
                        print(f"[Ingestion] Отмена во время OCR кадров ({idx}/{n})")
                        try: executor.shutdown(wait=False, cancel_futures=True)
                        except Exception: pass
                        raise IngestionCancelled(f"Cancelled at frame {idx}/{n}")
                    desc = future.result()
                    path, t = frame_list[idx]

                    # F3: если описание влезает в контекст эмбеддингов — оставляем одним чанком.
                    # Иначе — режем через сплиттер. Это сохраняет связность описания для
                    # средних картинок (2-4К символов) и не фрагментирует семантику.
                    full_text = f"Кадр {file_name} [{format_seconds(t)}]: {desc}"
                    if len(full_text) <= config.GGUF_CTX_EMBED_CHARS:
                        # Один цельный чанк
                        nodes.append(TextNode(text=full_text,
                                              metadata={"file_name":file_name, "image_path":path, "time":t}))
                    else:
                        # Только для очень длинных описаний — split
                        desc_nodes = splitter.get_nodes_from_documents([
                            TextNode(text=full_text,
                                     metadata={"file_name":file_name, "image_path":path, "time":t})
                        ])
                        nodes.extend(desc_nodes)

                    frame_data.append({"time":t, "image_path":path, "description":desc})
                    done = idx + 1
                    prog(65 + int(done/n*22) if n else 87, f"Описание: {done}/{n}")
            except IngestionCancelled:
                raise

        if shared_llm_url and not keep_vision_alive:
            unload_all_models(role="llm")

    metadata_json = {"file_name": file_name, "is_video": is_video, "transcript": transcript_data, "frames": frame_data}
    with open(os.path.join(os.path.dirname(file_path), f"{file_name}.json"), "w", encoding="utf-8") as f:
        json.dump(metadata_json, f, ensure_ascii=False, indent=2)
    return nodes

def _analyze_page_for_vision(page):
    """Чистый анализ одной страницы без побочных эффектов.
    Возвращает (text, has_real_graphics). Потокобезопасно: fitz.Page immutable из open документа."""
    text = page.get_text()
    images = page.get_images()
    drawings = page.get_drawings()

    has_real_graphics = False
    if len(images) > 0:
        has_real_graphics = True
    else:
        graphics_weight = 0
        horizontal_lines = 0
        vertical_lines = 0
        for d in drawings:
            items = d.get('items', [])
            if any(i[0] in ['c', 'q'] for i in items) or len(items) > 12:
                has_real_graphics = True
                break

            rect = d.get('rect')
            fill = d.get('fill')
            is_page_background = (
                len(items) == 1
                and items[0][0] == 're'
                and fill is not None
                and all(c >= 0.99 for c in fill)
                and rect is not None
                and (rect.x1 - rect.x0) > 200
                and (rect.y1 - rect.y0) > 200
            )
            if is_page_background:
                continue

            if rect is not None:
                w = rect.x1 - rect.x0
                h = rect.y1 - rect.y0
                if w > 30 and h < 3:
                    horizontal_lines += 1
                elif w < 3 and h > 30:
                    vertical_lines += 1
            graphics_weight += 1

        if not has_real_graphics:
            if graphics_weight > 8:
                has_real_graphics = True
            elif horizontal_lines >= 3 and vertical_lines >= 1:
                has_real_graphics = True
            elif horizontal_lines + vertical_lines >= 6:
                has_real_graphics = True

    return text, has_real_graphics


def process_pdf(file_path, images_dir, llm_settings=None, shared_llm_url=None, original_filename=None, progress_cb=None, cancel_check=None, keep_vision_alive=False):
    def _is_cancelled():
        return bool(cancel_check and cancel_check())

    nodes = []; file_name = original_filename or os.path.basename(file_path); doc = fitz.open(file_path)
    frame_data = []
    frame_list = []
    splitter = SentenceSplitter(chunk_size=1024, chunk_overlap=256)
    # Сплиттер для описаний (более крупный)
    v_splitter = SentenceSplitter(chunk_size=2048, chunk_overlap=128)

    # Параллельный разбор страниц: page.get_text() + get_drawings() GIL-free (PyMuPDF native).
    # Один ThreadPoolExecutor на весь PDF; cancel-check между батчами.
    # n_workers ограничен min(8, cpu_count, num_pages) чтобы не плодить потоки на маленьких PDF.
    total_pages = len(doc)
    n_workers = min(8, (os.cpu_count() or 4), total_pages)
    page_results = [None] * total_pages  # сохраняем порядок страниц

    if n_workers <= 1:
        # Тривиальный случай: последовательный проход
        for page_num in range(total_pages):
            if _is_cancelled():
                print(f"[Ingestion] Отмена на странице {page_num+1}/{total_pages}")
                raise IngestionCancelled(f"Cancelled at page {page_num+1}")
            page = doc.load_page(page_num)
            text, has_real_graphics = _analyze_page_for_vision(page)
            page_results[page_num] = (text, has_real_graphics)
    else:
        with ThreadPoolExecutor(max_workers=n_workers) as ex:
            cancel_check_every = max(1, total_pages // (n_workers * 4))
            submitted = 0
            future_to_pn = {}
            for page_num in range(total_pages):
                if submitted % cancel_check_every == 0 and _is_cancelled():
                    print(f"[Ingestion] Отмена во время разбора страниц ({submitted}/{total_pages})")
                    ex.shutdown(wait=False, cancel_futures=True)
                    raise IngestionCancelled(f"Cancelled at page {submitted}")
                page = doc.load_page(page_num)
                future_to_pn[ex.submit(_analyze_page_for_vision, page)] = page_num
                submitted += 1
            for fut in as_completed(future_to_pn):
                pn = future_to_pn[fut]
                try:
                    page_results[pn] = fut.result()
                except Exception as e:
                    logger.warning(f"Ошибка разбора страницы {pn+1}: {e}")
                    page_results[pn] = ("", False)

    # Параллельная фаза: создание нод + Pixmap-ов.
    # F-fix #22: PyMuPDF >= 1.24 (у нас 1.27.2.3) — get_pixmap() thread-safe.
    # Раньше делалось последовательно с комментарием "Pixmap НЕ потокобезопасен в
    # PyMuPDF < 1.24", но версия давно превысила 1.24, а комментарий остался.
    # Для 100-страничного PDF экономим до 5-7 сек на вводе-выводе.
    def _build_page_artifacts(page_num: int):
        text, has_real_graphics = page_results[page_num]
        local_nodes = []
        image_path = None
        if text and text.strip():
            local_nodes = splitter.get_nodes_from_documents([
                TextNode(text=text, metadata={"file_name":file_name, "page":page_num+1})
            ])
        if has_real_graphics:
            image_path = os.path.join(images_dir, f"p_{page_num+1}_{uuid.uuid4().hex[:6]}.png")
            doc.load_page(page_num).get_pixmap(dpi=150).save(image_path)
        return page_num, local_nodes, image_path

    if total_pages <= 1 or n_workers <= 1:
        # Последовательный путь (1 страница — оверхед потоков не оправдан)
        for page_num in range(total_pages):
            if _is_cancelled():
                print(f"[Ingestion] Отмена на странице {page_num+1}/{total_pages}")
                raise IngestionCancelled(f"Cancelled at page {page_num+1}")
            pn, local_nodes, image_path = _build_page_artifacts(page_num)
            nodes.extend(local_nodes)
            if image_path:
                frame_list.append({"page": pn+1, "path": image_path})
    else:
        with ThreadPoolExecutor(max_workers=n_workers) as ex:
            futures = [ex.submit(_build_page_artifacts, pn) for pn in range(total_pages)]
            artifacts = [None] * total_pages  # сохраняем порядок
            for fut in as_completed(futures):
                if _is_cancelled():
                    ex.shutdown(wait=False, cancel_futures=True)
                    raise IngestionCancelled("Cancelled during artifacts build")
                pn, local_nodes, image_path = fut.result()
                artifacts[pn] = (local_nodes, image_path)
            for pn, (local_nodes, image_path) in enumerate(artifacts):
                nodes.extend(local_nodes)
                if image_path:
                    frame_list.append({"page": pn+1, "path": image_path})

    if frame_list:
        # Ленивый запуск сервера
        if shared_llm_url is None:
            shared_llm_url = get_vision_url(llm_settings)
        
        if shared_llm_url:
            v_conc = int(llm_settings.get("vision_concurrency") or config.VISION_CONCURRENCY)
            n = len(frame_list)
            if progress_cb: progress_cb(65, f"Анализ {n} страниц PDF ({'параллельно' if v_conc > 1 else 'последовательно'})...")
            
            with ThreadPoolExecutor(max_workers=v_conc) as executor:
                futures = {executor.submit(describe_image_with_lmstudio, f["path"], llm_settings, shared_llm_url): f for f in frame_list}

                done_count = 0
                try:
                    for future in as_completed(futures):
                        if _is_cancelled():
                            print(f"[Ingestion] Отмена во время OCR ({done_count}/{n} готово)")
                            for f in futures.values():
                                try: executor.shutdown(wait=False, cancel_futures=True)
                                except Exception: pass
                            raise IngestionCancelled(f"Cancelled during OCR ({done_count}/{n})")
                        frame_info = futures[future]
                        desc = future.result()
                        done_count += 1

                        if desc and "Изображение без описания" not in desc:
                            # F3: если описание влезает в ctx эмбеддингов — оставляем одним чанком.
                            # Типичное vision-описание 1.5-4К символов; раньше v_splitter резал
                            # описание таблицы/схемы пополам, теряя связность.
                            full_text = f"Изображение PDF {file_name} стр {frame_info['page']}: {desc}"
                            if len(full_text) <= config.GGUF_CTX_EMBED_CHARS:
                                # Один цельный чанк — embedding получает полную семантику
                                nodes.append(TextNode(
                                    text=full_text,
                                    metadata={"file_name":file_name, "image_path":frame_info["path"], "page":frame_info["page"]}
                                ))
                            else:
                                # Только очень длинные описания (>4К) — split
                                desc_nodes = v_splitter.get_nodes_from_documents([
                                    TextNode(
                                        text=full_text,
                                        metadata={"file_name":file_name, "image_path":frame_info["path"], "page":frame_info["page"]}
                                    )
                                ])
                                nodes.extend(desc_nodes)

                            frame_data.append({
                                "page": frame_info["page"],
                                "image_path": frame_info["path"],
                                "description": desc
                            })
                        else:
                            try: os.remove(frame_info["path"])
                            except Exception: pass

                        if progress_cb: progress_cb(65 + int(done_count/n*25), f"Описание PDF: {done_count}/{n}")
                except IngestionCancelled:
                    raise

        if shared_llm_url and not keep_vision_alive:
            unload_all_models(role="llm")

    # Сохраняем метаданные для PDF (как для видео), чтобы фронтенд мог показать картинки
    if frame_data:
        # Сортируем по номеру страницы
        frame_data.sort(key=lambda x: x["page"])
        metadata_json = {"file_name": file_name, "is_video": False, "transcript": [], "frames": frame_data}
        with open(os.path.join(os.path.dirname(file_path), f"{file_name}.json"), "w", encoding="utf-8") as f:
            json.dump(metadata_json, f, ensure_ascii=False, indent=2)
            
    return nodes

def process_pptx(file_path, images_dir, llm_settings=None, shared_llm_url=None, progress_cb=None, cancel_check=None, keep_vision_alive=False):
    nodes = []; file_name = os.path.basename(file_path); pdf_path = os.path.splitext(file_path)[0] + ".pdf"
    import win32com.client, pythoncom
    app = None; deck = None
    try:
        pythoncom.CoInitialize()
        app = win32com.client.Dispatch("Powerpoint.Application")
        deck = app.Presentations.Open(os.path.abspath(file_path), WithWindow=False)
        deck.SaveAs(os.path.abspath(pdf_path), 32)
        # После конвертации удаляем оригинал и работаем с PDF
        if os.path.exists(pdf_path):
            if os.path.exists(file_path): os.remove(file_path)
            # Метаданные сохраняем уже для нового PDF
            nodes = process_pdf(pdf_path, images_dir, llm_settings, shared_llm_url, original_filename=os.path.basename(pdf_path), progress_cb=progress_cb, cancel_check=cancel_check, keep_vision_alive=keep_vision_alive)
        else:
            raise Exception("PDF conversion failed")
    except IngestionCancelled:
        raise
    except Exception as e:
        logger.warning(f"COM-конвертация PPTX не удалась, резерв через python-pptx: {e}")
        prs = Presentation(file_path)
        for i, slide in enumerate(prs.slides):
            nodes.append(TextNode(text="\n".join([sh.text for sh in slide.shapes if hasattr(sh, "text")]), metadata={"file_name":file_name, "page":i+1}))
    finally:
        # F-fix #14: гарантированно закрыть COM-объекты, иначе процесс PowerPoint
        # остаётся висеть с залоченным файлом. try/finally выполняется даже при exception.
        if deck is not None:
            try: deck.Close()
            except Exception: pass
        if app is not None:
            try: app.Quit()
            except Exception: pass
        try: pythoncom.CoUninitialize()
        except Exception: pass
    return nodes

def process_docx(file_path, images_dir, llm_settings=None, shared_llm_url=None, progress_cb=None, cancel_check=None, keep_vision_alive=False):
    nodes = []; file_name = os.path.basename(file_path); pdf_path = os.path.splitext(file_path)[0] + ".pdf"
    import win32com.client, pythoncom
    app = None; doc = None
    try:
        pythoncom.CoInitialize()
        app = win32com.client.Dispatch("Word.Application")
        doc = app.Documents.Open(os.path.abspath(file_path))
        doc.SaveAs(os.path.abspath(pdf_path), 17)
        # После конвертации удаляем оригинал и работаем с PDF
        if os.path.exists(pdf_path):
            if os.path.exists(file_path): os.remove(file_path)
            # Метаданные сохраняем уже для нового PDF
            nodes = process_pdf(pdf_path, images_dir, llm_settings, shared_llm_url, original_filename=os.path.basename(pdf_path), progress_cb=progress_cb, cancel_check=cancel_check, keep_vision_alive=keep_vision_alive)
        else:
            raise Exception("PDF conversion failed")
    except IngestionCancelled:
        raise
    except Exception as e:
        logger.warning(f"COM-конвертация DOCX не удалась, резерв через python-docx: {e}")
        import docx
        nodes.append(TextNode(text="\n".join([p.text for p in docx.Document(file_path).paragraphs]), metadata={"file_name":file_name}))
    finally:
        # F-fix #14: гарантированно закрыть COM-объекты Word (см. process_pptx).
        if doc is not None:
            try: doc.Close()
            except Exception: pass
        if app is not None:
            try: app.Quit()
            except Exception: pass
        try: pythoncom.CoUninitialize()
        except Exception: pass
    return nodes

def _safe_print(msg):
    """Print с защитой от cp1251 PowerShell (emoji/non-ASCII).

    F-fix #23: НЕ мутируем глобальный sys.stdout — это racy в
    ThreadPoolExecutor (один поток reconfigure()s, другой в этот момент
    print()-ит в старой кодировке → битый вывод). Вместо этого пишем
    напрямую в sys.stdout.buffer в utf-8 (он не зависит от text-stream
    кодировки sys.stdout).
    """
    try:
        print(msg)
    except UnicodeEncodeError:
        try:
            sys.stdout.buffer.write((str(msg) + "\n").encode("utf-8", errors="replace"))
            sys.stdout.buffer.flush()
        except Exception:
            print(msg.encode('ascii', errors='replace').decode('ascii'))


def ensure_720p_video(file_path, prog_cb=None, cancel_check=None, notebook_id=None):
    ext = os.path.splitext(file_path)[1].lower()
    if ext not in ['.mp4', '.avi', '.mkv', '.mov']: return file_path
    from imageio_ffmpeg import get_ffmpeg_exe
    ffmpeg = get_ffmpeg_exe()
    def _is_cancelled():
        return bool(cancel_check and cancel_check())
    def get_duration(path):
        # FIX: добавлен таймаут 30с и явный flush. Раньше `ffmpeg -i path` без таймаута
        # мог зависнуть на минуты для .mkv с нестандартными контейнерами.
        # 30с — за глаза для парсинга заголовка; если не успел — возвращаем 0
        # (caller корректно обрабатывает fallback).
        try:
            cmd = [ffmpeg, "-hide_banner", "-i", path]
            res = subprocess.run(cmd, capture_output=True, timeout=30)
            # F-fix #31: capture_output=True теоретически гарантирует bytes, но
            # на Windows + специфические ffmpeg-ошибки могут дать res.stderr=None.
            stderr = (res.stderr or b"").decode('utf-8', errors='ignore')
            for line in stderr.split("\n"):
                if "Duration" in line:
                    time_str = line.split("Duration: ")[1].split(",")[0]
                    h, m, s = time_str.split(":"); return float(h)*3600 + float(m)*60 + float(s)
        except subprocess.TimeoutExpired:
            _safe_print(f"[ensure_720p_video] WARNING get_duration timeout для {os.path.basename(path)} (30с)")
        except Exception: pass
        return 0
    _safe_print(f"[ensure_720p_video] Начало: {os.path.basename(file_path)}")
    duration = get_duration(file_path); temp_final = file_path + ".720p.mp4"
    _safe_print(f"[ensure_720p_video] Длительность: {duration:.1f}с ({format_seconds(duration)})")
    # FIX: если get_duration упал (вернул 0) — идём в TURBO-режим (быстрее для длинных),
    # а не в single-encode (который медленнее для 1+ч видео). Раньше duration=0
    # попадал в ветку duration<120 → single encode → 30-60 мин вместо 5-10 мин.
    if duration == 0:
        _safe_print(f"[ensure_720p_video] WARNING Длительность неизвестна -> турбо-режим по умолчанию")
        use_turbo = True
    else:
        use_turbo = duration >= 120
    if not use_turbo:
        if _is_cancelled(): raise IngestionCancelled("Cancelled before 720p encode")
        if prog_cb: prog_cb(5, "Оптимизация видео (GPU)...")
        # FIX: stderr в DEVNULL — раньше PIPE заполнялся за 30-60с и ffmpeg блокировался
        cmd = [ffmpeg, "-y", "-hide_banner", "-hwaccel", "cuda", "-hwaccel_output_format", "cuda", "-i", file_path, "-vf", "scale_cuda=1280:720:format=yuv420p", "-c:v", "hevc_nvenc", "-preset", "p1", "-rc", "constqp", "-qp", "30", "-pix_fmt", "yuv420p", "-tag:v", "hvc1", "-c:a", "aac", "-b:a", "128k", "-movflags", "+faststart", temp_final]
        proc = subprocess.Popen(cmd, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
        if notebook_id is not None: register_subprocess(notebook_id, proc)
        try:
            proc.wait()
        finally:
            if notebook_id is not None: unregister_subprocess(notebook_id, proc)
        if _is_cancelled():
            try: os.remove(temp_final)
            except Exception: pass
            raise IngestionCancelled("Cancelled during 720p encode")
    else:
        if _is_cancelled(): raise IngestionCancelled("Cancelled before turbo encode")
        if prog_cb: prog_cb(5, "Турбо-оптимизация (Параллельный GPU)...")
        num_workers = 4; seg_len = duration / num_workers; temp_dir = file_path + "_parts"
        os.makedirs(temp_dir, exist_ok=True)
        def encode_seg(idx):
            if _is_cancelled(): return None
            out_part = os.path.join(temp_dir, f"part_{idx}.mp4")
            # FIX: stderr в DEVNULL (раньше PIPE заполнялся → ffmpeg блокировался →
            # все 4 параллельных сегмента зависали одновременно через 30-60с).
            # stdout с -progress pipe:1 → читаем построчно → real-time progress в UI.
            cmd = [ffmpeg, "-y", "-hide_banner", "-hwaccel", "cuda", "-hwaccel_output_format", "cuda",
                   "-ss", str(idx * seg_len), "-t", str(seg_len), "-i", file_path,
                   "-vf", "scale_cuda=1280:720:format=yuv420p",
                   "-c:v", "hevc_nvenc", "-preset", "p1", "-rc", "constqp", "-qp", "30",
                   "-progress", "pipe:1", "-an", out_part]
            proc = subprocess.Popen(cmd, stdout=subprocess.PIPE, stderr=subprocess.DEVNULL,
                                    text=True, bufsize=1)
            if notebook_id is not None: register_subprocess(notebook_id, proc)
            try:
                last_pct = -1
                for line in proc.stdout:
                    if _is_cancelled():
                        proc.kill(); return None
                    line = line.strip()
                    if line.startswith("out_time_ms="):
                        try:
                            time_ms = int(line.split("=", 1)[1])
                            pct = min(99, int(time_ms / 10000 / seg_len))
                            if pct > last_pct and pct % 10 == 0:
                                last_pct = pct
                                if prog_cb:
                                    overall = 5 + ((idx + pct / 100.0) / num_workers) * 3
                                    prog_cb(overall, f"Сегмент {idx+1}/{num_workers}: {pct}%")
                        except (ValueError, ZeroDivisionError):
                            pass
                proc.wait()
            finally:
                if notebook_id is not None: unregister_subprocess(notebook_id, proc)
            return out_part
        with concurrent.futures.ThreadPoolExecutor(max_workers=num_workers) as ex:
            parts = []
            for p in ex.map(encode_seg, range(num_workers)):
                if _is_cancelled():
                    _safe_print(f"[Ingestion] Отмена во время турбо-кодирования видео (получено {len(parts)}/{num_workers} сегментов)")
                    ex.shutdown(wait=False, cancel_futures=True)
                    raise IngestionCancelled("Cancelled during turbo encode")
                parts.append(p)
        if _is_cancelled(): raise IngestionCancelled("Cancelled after turbo encode")
        if prog_cb: prog_cb(8, "Сборка сегментов...")
        list_path = os.path.join(temp_dir, "list.txt")
        with open(list_path, "w") as f:
            for p in parts: f.write(f"file '{os.path.abspath(p)}'\n")
        # FIX: merge тоже с DEVNULL для stderr
        merge_cmd = [ffmpeg, "-y", "-hide_banner", "-f", "concat", "-safe", "0", "-i", list_path, "-i", file_path, "-map", "0:v", "-map", "1:a?", "-c", "copy", "-movflags", "+faststart", temp_final]
        merge_proc = subprocess.Popen(merge_cmd, stdout=subprocess.DEVNULL, stderr=subprocess.DEVNULL)
        if notebook_id is not None: register_subprocess(notebook_id, merge_proc)
        try:
            merge_proc.wait()
        finally:
            if notebook_id is not None: unregister_subprocess(notebook_id, merge_proc)
        if _is_cancelled():
            try: os.remove(temp_final)
            except Exception: pass
            raise IngestionCancelled("Cancelled after video merge")
        try: shutil.rmtree(temp_dir)
        except Exception: pass
    if os.path.exists(temp_final) and os.path.getsize(temp_final) > 1000:
        if os.path.exists(file_path): os.remove(file_path)
        new_path = os.path.splitext(file_path)[0] + ".mp4"
        if os.path.exists(new_path) and new_path != temp_final: os.remove(new_path)
        os.rename(temp_final, new_path); file_path = new_path
        if prog_cb: prog_cb(9, "Видео оптимизировано (Турбо)")
    return file_path

def ensure_mp3_audio(file_path, prog_cb=None):
    temp_path = file_path.rsplit('.', 1)[0] + ".mp3"
    from imageio_ffmpeg import get_ffmpeg_exe
    cmd = [get_ffmpeg_exe(), "-y", "-i", file_path, "-acodec", "libmp3lame", "-ab", "128k", temp_path]
    subprocess.run(cmd, capture_output=True)
    if os.path.exists(temp_path): os.remove(file_path); return temp_path
    return file_path

def ingest_file(file_path, notebook_id, progress_cb=None, llm_settings=None, cancel_check=None, keep_vision_alive=False, keep_whisper_alive=False):
    def _is_cancelled():
        return bool(cancel_check and cancel_check())

    ext = os.path.splitext(file_path)[1].lower()
    paths = config.get_notebook_paths(notebook_id); images_dir = paths["images"]
    os.makedirs(images_dir, exist_ok=True)
    if _is_cancelled(): raise IngestionCancelled("Cancelled before media conversion")
    if ext in ['.mp4', '.avi', '.mkv', '.mov']: file_path = ensure_720p_video(file_path, progress_cb, cancel_check=cancel_check, notebook_id=notebook_id)
    elif ext in ['.mp3', '.wav', '.m4a']: file_path = ensure_mp3_audio(file_path, progress_cb); ext = ".mp3"
    if _is_cancelled(): raise IngestionCancelled("Cancelled after media conversion")
    if ext in ['.mp4', '.avi', '.mkv', '.mov', '.mp3']: return process_audio_video(file_path, images_dir, ext != ".mp3", progress_cb, llm_settings, cancel_check=cancel_check, notebook_id=notebook_id, keep_vision_alive=keep_vision_alive, keep_whisper_alive=keep_whisper_alive)
    # Больше не запускаем Vision-сервер заранее.
    # Он запустится лениво (lazy-load) только если внутри PDF/PPTX/DOCX обнаружится реальное изображение.
    shared_llm_url = None

    try:
        if ext == '.pdf': nodes = process_pdf(file_path, images_dir, llm_settings, shared_llm_url, progress_cb=progress_cb, cancel_check=cancel_check, keep_vision_alive=keep_vision_alive)
        elif ext == '.pptx': nodes = process_pptx(file_path, images_dir, llm_settings, shared_llm_url, progress_cb=progress_cb, cancel_check=cancel_check, keep_vision_alive=keep_vision_alive)
        elif ext == '.docx': nodes = process_docx(file_path, images_dir, llm_settings, shared_llm_url, progress_cb=progress_cb, cancel_check=cancel_check, keep_vision_alive=keep_vision_alive)
        else:
            try:
                with open(file_path, 'r', encoding='utf-8') as f: text = f.read()
            except UnicodeDecodeError:
                with open(file_path, 'r', encoding='cp1251') as f: text = f.read()
                
            doc = TextNode(text=text, metadata={"file_name":os.path.basename(file_path)})
            nodes = SentenceSplitter(chunk_size=1024, chunk_overlap=256).get_nodes_from_documents([doc])
    finally:
        # Мы НЕ выгружаем модели в finally, чтобы они оставались в памяти для следующего файла в батче.
        # Модели будут выгружены автоматически, если потребуется память для другого типа моделей.
        pass
    return nodes
