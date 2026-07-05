"""Обработка текстовых документов: PDF, PPTX, DOCX."""

import asyncio
import logging
import os
import uuid

import aiofiles
import fitz
import orjson
from llama_index.core.schema import TextNode

import config
from routers.shared import get_async_http
from src.gguf.server import unload_all_models
from src.ingestion.splitter import _get_splitter
from src.ingestion.utils import IngestionCancelled
from src.ingestion.vision import describe_image_with_lmstudio, get_vision_url

logger = logging.getLogger(__name__)


<<<<<<< HEAD
# Извлечение встроенных изображений со страницы PDF
# (только когда surya_mode=disabled — иначе surya layout делает это лучше)


async def _analyze_and_build_page(
    page_num, doc, images_dir, file_name, splitter, surya_mode="disabled"
):

    if surya_mode in ("layout_only", "full"):
        return page_num, []

    def _sync_build():
        page = doc.load_page(page_num)
        images = page.get_images()

        image_paths = []
        seen_xrefs = set()
        for img_info in images:
            xref = img_info[0]
            if xref in seen_xrefs:
                continue
            seen_xrefs.add(xref)
            try:
                rects = page.get_image_rects(xref)
                if not rects:
                    continue
                base_image = page.parent.extract_image(xref)
                if not base_image or not base_image.get("image"):
                    continue
                ext = base_image.get("ext", "png")
                if ext == "jpeg":
                    ext = "jpg"
                for rect in rects:
                    w, h = rect.width, rect.height
                    if w < 50 or h < 50:
                        continue
                    try:
                        clip = fitz.Rect(rect)
                        pix = page.get_pixmap(clip=clip, dpi=150)
                        img_name = f"p_{page_num + 1}_x{xref}_{uuid.uuid4().hex[:4]}.{ext}"
                        img_path = os.path.join(images_dir, img_name)
                        pix.save(img_path)
                        image_paths.append(img_path)
                    except Exception:
                        continue
            except Exception:
                continue

        return page_num, image_paths

    return await asyncio.to_thread(_sync_build)


# Основной конвейер PDF: извлечение текста, батчевая параллельная обработка, Vision-анализ
=======
# Основной конвейер PDF: извлечение текста (pymupdf4llm), изображения (surya layout), Vision-анализ
>>>>>>> 62b6f39 (refactor(async): замена asyncio.to_thread(os.path.exists) на aiofiles.os.path.exists)
async def process_pdf(
    file_path,
    images_dir,
    llm_settings=None,
    shared_llm_url=None,
    original_filename=None,
    progress_cb=None,
    cancel_check=None,
    keep_vision_alive=False,
):

    def _is_cancelled():
        return bool(cancel_check and cancel_check())

    nodes = []
    file_name = original_filename or os.path.basename(file_path)
    doc = await asyncio.to_thread(fitz.open, file_path)
    frame_data = []
    frame_list = []
    results = []
    splitter = _get_splitter()
    total_pages = len(doc)

    try:
        import pymupdf4llm

        def _extract_markdown():
            return pymupdf4llm.to_markdown(file_path, page_chunks=True, write_images=False)

        if progress_cb:
            progress_cb(10, "Извлечение текста (pymupdf4llm)...")
        md_chunks = await asyncio.to_thread(_extract_markdown)

<<<<<<< HEAD
        def _clean_markdown(text: str) -> str:
            import re

            text = re.sub(r"^(#{1,6})\s*\*\*(.+?)\*\*\s*$", r"\1 \2", text, flags=re.MULTILINE)
            text = re.sub(r"\*\*==>.+?intentionally omitted.+?<==\*\*", "", text)
            text = re.sub(r"\*\*(.+?\.{3,}\d+)\*\*", r"\1", text)
            text = re.sub(r"\n\s*\d{1,3}\s*\n", "\n", text)
            text = text.replace("\\n", "\n")
            text = re.sub(r"\n{3,}", "\n\n", text)
            return text
=======
        # Оцениваем качество извлечения
        total_chars = sum(len(c.get("text", "")) for c in md_chunks)
        avg_chars = total_chars / max(total_pages, 1)
        logger.info(
            f"[Ingestion] pymupdf4llm: {len(md_chunks)} чанков, {total_chars} симв., среднее {avg_chars:.0f}/стр"
        )

        if avg_chars < 100:
            # PDF-скан, мало текста → нужен Surya OCR
            logger.info("[Ingestion] Мало текста (< 100/стр) → Surya OCR")
            use_surya_ocr = True
        else:
            # Хороший текстовый PDF
            import re

            # ── Очистка markdown от артефактов pymupdf ──
            def _clean_markdown(text):
                text = re.sub(r"^(#{1,6})\s*\*\*(.+?)\*\*\s*$", r"\1 \2", text, flags=re.MULTILINE)
                text = re.sub(
                    r"\*{0,2}={1,2}> picture \[\d+ x \d+\] intentionally omitted <={1,2}\*{0,2}",
                    "",
                    text,
                )
                text = re.sub(
                    r"----- Start of picture text -----.*?----- End of picture text -----",
                    "",
                    text,
                    flags=re.DOTALL,
                )
                text = re.sub(r"\*\*(.+?\.{3,}\d+)\*\*", r"\1", text)
                text = re.sub(r"\n\s*\d{1,3}\s*\n", "\n", text)
                text = text.replace("\\n", "\n")
                text = re.sub(r"\n{3,}", "\n\n", text)
                text = text.strip()
                return text
>>>>>>> 62b6f39 (refactor(async): замена asyncio.to_thread(os.path.exists) на aiofiles.os.path.exists)

        for chunk in md_chunks:
            page_num = chunk.get("metadata", {}).get("page_number", 0)
            md_text = chunk.get("text", "")
            if md_text and md_text.strip():
                md_text = _clean_markdown(md_text)
                nodes.extend(
                    splitter.get_nodes_from_documents(
                        [
                            TextNode(
                                text=md_text,
                                metadata={
                                    "file_name": file_name,
                                    "page": page_num,
                                },
                            )
                        ]
                    )
                )
        logger.info(f"[Ingestion] pymupdf4llm: {len(md_chunks)} чанков, {len(nodes)} узлов")

        # ── Сохраняем извлечённый текст на диск для мгновенного доступа ──
        try:
            extracted_parts = []
            for chunk in md_chunks:
<<<<<<< HEAD
                t = chunk.get("text", "")
                if t and t.strip():
                    extracted_parts.append(t)
            if extracted_parts:
                extracted_path = os.path.join(os.path.dirname(file_path), f"{file_name}.extracted.md")
                async with aiofiles.open(extracted_path, "w", encoding="utf-8") as ef:
                    await ef.write("\n\n---\n\n".join(extracted_parts))
        except Exception as e:
            logger.warning(f"[Ingestion] Не удалось сохранить .extracted.md: {e}")
=======
                page_num = chunk.get("metadata", {}).get("page_number", 0)
                md_text = chunk.get("text", "")
                if md_text and md_text.strip():
                    md_text = _clean_markdown(md_text)

                    # Удаляем pymupdf image placeholders — surya layout извлекает изображения
                    md_text = re.sub(
                        r"=> picture \[\d+ x \d+\] intentionally omitted <=.*?(?=----- Start|$)",
                        "",
                        md_text,
                        flags=re.DOTALL,
                    )
                    md_text = re.sub(
                        r"----- Start of picture text -----.*?----- End of picture text -----",
                        "",
                        md_text,
                        flags=re.DOTALL,
                    )
                    md_text = re.sub(r"\n{3,}", "\n\n", md_text).strip()

                    if md_text and md_text.strip():
                        # Добавляем номер страницы
                        page_header = f"## Стр. {page_num}\n\n"
                        nodes.extend(
                            splitter.get_nodes_from_documents(
                                [
                                    TextNode(
                                        text=page_header + md_text,
                                        metadata={"file_name": file_name, "page": page_num},
                                    )
                                ]
                            )
                        )

            logger.info(f"[Ingestion] pymupdf4llm: {len(nodes)} узлов")

>>>>>>> 62b6f39 (refactor(async): замена asyncio.to_thread(os.path.exists) на aiofiles.os.path.exists)
    except ImportError:
        logger.warning("[Ingestion] pymupdf4llm не установлен — fallback на page.get_text()")
        for page_num in range(total_pages):
            page = doc.load_page(page_num)
            text = page.get_text()
            if text and text.strip():
                nodes.extend(
                    splitter.get_nodes_from_documents(
                        [
                            TextNode(
                                text=text,
                                metadata={
                                    "file_name": file_name,
                                    "page": page_num + 1,
                                },
                            )
                        ]
                    )
                )

    surya_mode = getattr(config, "SURYA_MODE", "disabled")
    n_workers = min(8, (os.cpu_count() or 4), total_pages)

    try:
        BATCH_SIZE = 16
        if n_workers <= 1:
            for page_num in range(total_pages):
                if _is_cancelled():
                    raise IngestionCancelled(f"Cancelled at page {page_num + 1}")
                pn, image_paths = await _analyze_and_build_page(
                    page_num, doc, images_dir, file_name, splitter, surya_mode
                )
                for img_p in image_paths or []:
                    frame_list.append({"page": pn + 1, "path": img_p})
        else:
            for batch_start in range(0, total_pages, BATCH_SIZE):
                if _is_cancelled():
                    raise IngestionCancelled(f"Cancelled at page {batch_start + 1}")
                batch_end = min(batch_start + BATCH_SIZE, total_pages)

                async def _process_batch_pages(batch_start, batch_end):
                    tasks = []
                    for page_num in range(batch_start, batch_end):
                        tasks.append(
                            _analyze_and_build_page(
                                page_num, doc, images_dir, file_name, splitter, surya_mode
                            )
                        )
                    return await asyncio.gather(*tasks)

                page_results = await _process_batch_pages(batch_start, batch_end)
                for page_num_result, image_paths in page_results:
                    for img_p in image_paths or []:
                        frame_list.append({"page": page_num_result + 1, "path": img_p})

        surya_mode = getattr(config, "SURYA_MODE", "disabled")
        if surya_mode in ("layout_only", "full") and not _is_cancelled():
            surya_frame_list = await _surya_layout_pass(
                doc,
                file_name,
                images_dir,
                splitter,
                nodes,
                llm_settings,
                shared_llm_url,
                progress_cb,
                cancel_check,
                surya_mode,
                frame_data,
                keep_vision_alive,
            )
            frame_list.extend(surya_frame_list)

        if frame_list:
            if shared_llm_url is None:
                shared_llm_url = await get_vision_url(llm_settings)
            if shared_llm_url:
                v_conc = int(
                    (llm_settings or {}).get("vision_concurrency") or config.VISION_CONCURRENCY
                )
                n = len(frame_list)
                if progress_cb:
                    progress_cb(
                        65,
                        f"Анализ {n} страниц PDF ({'параллельно' if v_conc > 1 else 'последовательно'})...",
                    )

                sem = asyncio.Semaphore(v_conc)
                done_count = 0
                results = []
                VISION_BATCH_SIZE = 100

                async def _describe_frame(frame_info):
                    nonlocal done_count
                    if _is_cancelled():
                        raise IngestionCancelled(f"Cancelled during OCR ({done_count}/{n})")
                    async with sem:
                        desc = await describe_image_with_lmstudio(
                            frame_info["path"],
                            llm_settings,
                            shared_llm_url,
                            cancel_check=cancel_check,
                        )
                        done_count += 1
                        results.append((frame_info, desc))
                        if progress_cb:
                            progress_cb(
                                65 + int(done_count / n * 25), f"Описание PDF: {done_count}/{n}"
                            )

                for batch_start in range(0, n, VISION_BATCH_SIZE):
                    if _is_cancelled():
                        raise IngestionCancelled(f"Cancelled at batch {batch_start}")
                    batch_end = min(batch_start + VISION_BATCH_SIZE, n)
                    batch = frame_list[batch_start:batch_end]
                    await asyncio.gather(*[_describe_frame(f) for f in batch])

                    if batch_end < n:
                        try:
                            http = await get_async_http()
                            resp = await http.get(f"{shared_llm_url}/health", timeout=2)
                            resp.raise_for_status()
                        except Exception:
                            logger.warning("[Vision] Сервер unhealthy, перезапуск...")
                            await unload_all_models(role="vision")
                            shared_llm_url = await get_vision_url(llm_settings)
                            if not shared_llm_url:
                                logger.warning("[Vision] Не удалось перезапустить vision-сервер")
                                break

                results.sort(key=lambda x: x[0]["page"])
                for frame_info, desc in results:
                    if desc and "Изображение без описания" not in desc:
<<<<<<< HEAD
                        full_text = f"Изображение PDF {file_name} стр {frame_info['page']}: {desc}"
                        if len(full_text) <= config.GGUF_CTX_EMBED_CHARS:
                            nodes.append(
                                TextNode(
                                    text=full_text,
                                    metadata={
                                        "file_name": file_name,
                                        "image_path": frame_info["path"],
                                        "page": frame_info["page"],
                                    },
                                )
                            )
                        else:
                            desc_nodes = splitter.get_nodes_from_documents(
                                [
                                    TextNode(
                                        text=full_text,
                                        metadata={
                                            "file_name": file_name,
                                            "image_path": frame_info["path"],
                                            "page": frame_info["page"],
                                        },
                                    )
                                ]
                            )
                            nodes.extend(desc_nodes)
                        frame_data.append(
                            {
                                "page": frame_info["page"],
                                "image_path": frame_info["path"],
                                "description": desc,
                            }
                        )
                    else:
                        try:
                            await aiofiles.os.remove(frame_info["path"])
                        except Exception:
                            pass

            results.clear()
=======
                        pg = frame_info["page"]
                        if pg not in page_descs:
                            page_descs[pg] = []
                        page_descs[pg].append(
                            {
                                "text": desc,
                                "image_path": frame_info["path"],
                            }
                        )

                # Добавляем описания как ОТДЕЛЬНЫЕ ноды (не обрезаются сплиттером)
                for pg, descs in list(page_descs.items()):
                    for d in descs:
                        desc_text = f"--- Изображение (стр. {pg}) ---\n{d['text']}\n---"
                        nodes.append(
                            TextNode(
                                text=desc_text,
                                metadata={
                                    "file_name": file_name,
                                    "page": pg,
                                    "image_path": d["image_path"],
                                },
                            )
                        )
                        frame_data.append(
                            {"page": pg, "image_path": d["image_path"], "description": d["text"]}
                        )

        if frame_data:
>>>>>>> 62b6f39 (refactor(async): замена asyncio.to_thread(os.path.exists) на aiofiles.os.path.exists)
            import gc

            gc.collect()

            if shared_llm_url and not keep_vision_alive:
                await unload_all_models(role="vision")

        if frame_data:
            frame_data.sort(key=lambda x: x["page"])
            metadata_json = {
                "file_name": file_name,
                "is_video": False,
                "transcript": [],
                "frames": frame_data,
            }

            metadata_path = os.path.join(os.path.dirname(file_path), f"{file_name}.json")
            async with aiofiles.open(metadata_path, "w", encoding="utf-8") as f:
                await f.write(orjson.dumps(metadata_json, option=orjson.OPT_INDENT_2).decode())
    finally:
        await asyncio.to_thread(doc.close)
    return nodes


# Surya layout pass: определение Diagram/Equation/Table regions + OCR
async def _surya_layout_pass(
    doc,
    file_name,
    images_dir,
    splitter,
    nodes,
    llm_settings,
    shared_llm_url,
    progress_cb,
    cancel_check,
    surya_mode,
    frame_data,
    keep_vision_alive,
):
    """Surya layout detection + OCR + Vision LLM для Diagram/Equation regions."""
    from src.ingestion.surya_layout import (
        detect_layout,
        extract_regions,
        ocr_text,
    )
    from src.ingestion.surya_layout import (
        shutdown as surya_shutdown,
    )

    def _is_cancelled():
        return bool(cancel_check and cancel_check())

    total_pages = len(doc)
    frame_list = []

    # Конвертируем страницы PDF в PIL images для surya
    if progress_cb:
        progress_cb(60, "Surya: определение layout...")

    def _render_pages():
        pil_images = []
        for page_num in range(total_pages):
            page = doc.load_page(page_num)
            pix = page.get_pixmap(dpi=150)
            from PIL import Image as _Image

            img = _Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
            pil_images.append(img)
        return pil_images

    pil_images = await asyncio.to_thread(_render_pages)

    # Запускаем layout detection
    try:
        layout_results = await asyncio.to_thread(detect_layout, pil_images)
    except Exception as e:
        logger.warning(f"[Surya] Ошибка layout detection: {e}")
        return frame_list

    if not layout_results:
        return frame_list

    if surya_mode == "full":
        try:
            if progress_cb:
                progress_cb(65, "Surya: OCR текста...")
            ocr_results = await asyncio.to_thread(ocr_text, pil_images, layout_results)
            if ocr_results:
                old_nodes = nodes.copy()
                nodes.clear()
                for page_data in ocr_results:
                    html_text = page_data.get("html", "")
                    if html_text and html_text.strip():
                        nodes.extend(
                            splitter.get_nodes_from_documents(
                                [
                                    TextNode(
                                        text=html_text,
                                        metadata={
                                            "file_name": file_name,
                                            "page": page_data["page"],
                                        },
                                    )
                                ]
                            )
                        )
                logger.info(f"[Surya] OCR: заменено {len(old_nodes)} -> {len(nodes)} узлов")
        except Exception as e:
            logger.warning(f"[Surya] Ошибка OCR: {e}")

    regions = {"Diagram", "Equation", "Table"}
    extracted = extract_regions(pil_images, layout_results, regions)

    n_regions = sum(len(r) for r in extracted.values())
    if n_regions == 0:
        logger.info("[Surya] Diagram/Equation/Table regions не найдены")
        surya_shutdown()
        return frame_list

    logger.info(f"[Surya] Найдено {n_regions} regions для описания")

    for page_idx, page_regions in extracted.items():
        for region in page_regions:
            if _is_cancelled():
                break
            try:
                label = region["label"]
                img = region["image"]
                img_name = f"surya_{label.lower()}_p{page_idx + 1}_{uuid.uuid4().hex[:4]}.png"
                img_path = os.path.join(images_dir, img_name)
                await asyncio.to_thread(img.save, img_path)
                frame_list.append({"page": page_idx + 1, "path": img_path})
            except Exception:
<<<<<<< HEAD
                logger.debug(f"[surya_layout] Не удалось сохранить region на странице {page_idx + 1}")
=======
                logger.debug(
                    f"[surya_layout] Не удалось сохранить группу layout на странице {page_idx + 1}"
                )
>>>>>>> 62b6f39 (refactor(async): замена asyncio.to_thread(os.path.exists) на aiofiles.os.path.exists)

    surya_shutdown()
    return frame_list


# Конвертация Office-файлов в PDF: LibreOffice (приоритет), COM (резерв), текстовый fallback
def _find_soffice():
    import platform
    import shutil

    # ── env override ──
    env_path = os.getenv("LIBREOFFICE_PATH")
    if env_path and os.path.isfile(env_path):
        return env_path

    # ── shutil.which — работает на всех платформах ──
    found = shutil.which("soffice") or shutil.which("soffice.exe")
    if found:
        return found

    # ── Платформо-специфичные пути ──
    if platform.system() == "Windows":
        local = os.path.join(config.BASE_DIR, "libreoffice", "program", "soffice.exe")
        if os.path.isfile(local):
            return local
        for pf in ["Program Files", "Program Files (x86)"]:
            p = os.path.join("C:\\", pf, "LibreOffice", "program", "soffice.exe")
            if os.path.isfile(p):
                return p
    else:
        for p in ["/usr/bin/soffice", "/usr/lib/libreoffice/program/soffice"]:
            if os.path.isfile(p):
                return p

    return None


async def _convert_via_libreoffice(file_path):
    soffice = _find_soffice()
    if not soffice:
        raise FileNotFoundError("LibreOffice не найден. Установите или скачайте через setup.ps1")
    out_dir = os.path.dirname(file_path)
    cmd = [
        soffice,
        "--headless",
        "--norestore",
        "--convert-to",
        "pdf",
        "--outdir",
        out_dir,
        os.path.abspath(file_path),
    ]
    proc = await asyncio.create_subprocess_exec(
        *cmd,
        stdout=asyncio.subprocess.PIPE,
        stderr=asyncio.subprocess.PIPE,
    )
    _stdout, stderr = await asyncio.wait_for(proc.communicate(), timeout=120)
    pdf_path = os.path.splitext(file_path)[0] + ".pdf"
    if proc.returncode == 0 and await aiofiles.os.path.exists(pdf_path):
        logger.info(f"[DOCX] Сконвертировано в PDF через LibreOffice: {os.path.basename(pdf_path)}")
        try:
            await aiofiles.os.remove(file_path)
        except OSError:
            pass
        return pdf_path
    raise RuntimeError(
        f"LibreOffice конвертация не удалась (code={proc.returncode}): "
        f"{(stderr or b'').decode('utf-8', errors='replace')[:200]}"
    )


async def _convert_via_com(file_path, app_name, format_code):

    def _sync_com():
        import pythoncom
        import win32com.client

        pdf_path = os.path.splitext(file_path)[0] + ".pdf"
        app = None
        doc = None
        try:
            pythoncom.CoInitialize()
            app = win32com.client.Dispatch(app_name)
            if app_name == "Powerpoint.Application":
                doc = app.Presentations.Open(os.path.abspath(file_path), WithWindow=False)
            else:
                doc = app.Documents.Open(os.path.abspath(file_path))
            doc.SaveAs(os.path.abspath(pdf_path), format_code)
            if not os.path.exists(pdf_path):
                raise RuntimeError(f"COM {app_name} не создал PDF: {pdf_path}")
            if os.path.exists(file_path):
                os.remove(file_path)
            return pdf_path
        except IngestionCancelled:
            raise
        except Exception as e:
            logger.warning(f"COM-конвертация через {app_name} не удалась: {e}")
            raise
        finally:
            if doc is not None:
                try:
                    doc.Close()
                except Exception:
                    pass
            if app is not None:
                try:
                    app.Quit()
                except Exception:
                    pass
            try:
                pythoncom.CoUninitialize()
            except Exception:
                pass

    return await asyncio.to_thread(_sync_com)


async def _convert_office_to_pdf(file_path, app_name, format_code, textonly_fn, log_prefix):
    file_name = os.path.basename(file_path)
    try:
        pdf_path = await _convert_via_libreoffice(file_path)
    except FileNotFoundError:
        logger.info(f"{log_prefix} LibreOffice не найден, пробую COM-конвертацию...")
        try:
            pdf_path = await _convert_via_com(file_path, app_name, format_code)
        except IngestionCancelled:
            raise
        except Exception as e:
            logger.warning(f"{log_prefix} COM-конвертация тоже не удалась: {e}")
            return None, await textonly_fn(file_path, file_name)
    except IngestionCancelled:
        raise
    except Exception as e:
        logger.warning(f"{log_prefix} LibreOffice конвертация не удалась: {e}")
        try:
            pdf_path = await _convert_via_com(file_path, app_name, format_code)
        except IngestionCancelled:
            raise
        except Exception:
            return None, await textonly_fn(file_path, file_name)
    return os.path.basename(pdf_path), None


def _pptx_extract_textonly(file_path, file_name):
    from pptx import Presentation

    prs = Presentation(file_path)
    nodes = []
    for i, slide in enumerate(prs.slides):
        text = "\n".join([sh.text for sh in slide.shapes if hasattr(sh, "text")])
        if text.strip():
            nodes.append(TextNode(text=text, metadata={"file_name": file_name, "page": i + 1}))
    if nodes:
        logger.info(
            f"[PPTX] Резервный вариант: {len(nodes)} слайдов через python-pptx (без Vision)"
        )
    return nodes


def _docx_extract_textonly(file_path, file_name):
    import docx as _docx

    text = "\n".join([p.text for p in _docx.Document(file_path).paragraphs])
    if text.strip():
        logger.info("[DOCX] Fallback: текст извлечён через python-docx (без Vision)")
        return [TextNode(text=text, metadata={"file_name": file_name})]
    return []


async def _process_office_textonly(file_path, file_name, extract_fn, log_prefix):

    def _sync():
        try:
            return extract_fn(file_path, file_name)
        except Exception as e:
            logger.warning(f"{log_prefix} резервный вариант тоже не удался: {e}")
        return []

    return await asyncio.to_thread(_sync)


# Обработка PPTX/DOCX: конвертация в PDF + Vision или текстовый fallback
async def process_pptx(
    file_path,
    images_dir,
    llm_settings=None,
    shared_llm_url=None,
    progress_cb=None,
    cancel_check=None,
    keep_vision_alive=False,
):
    file_name = os.path.basename(file_path)
    pdf_name, fallback_nodes = await _convert_office_to_pdf(
        file_path,
        "Powerpoint.Application",
        32,
        lambda fp, fn: _process_office_textonly(fp, fn, _pptx_extract_textonly, "[PPTX]"),
        "[PPTX]",
    )
    if fallback_nodes is not None:
        return fallback_nodes

    return await process_pdf(
        os.path.join(os.path.dirname(file_path), pdf_name),
        images_dir,
        llm_settings,
        shared_llm_url,
        original_filename=pdf_name,
        progress_cb=progress_cb,
        cancel_check=cancel_check,
        keep_vision_alive=keep_vision_alive,
    )


async def process_docx(
    file_path,
    images_dir,
    llm_settings=None,
    shared_llm_url=None,
    progress_cb=None,
    cancel_check=None,
    keep_vision_alive=False,
):
    file_name = os.path.basename(file_path)
    pdf_name, fallback_nodes = await _convert_office_to_pdf(
        file_path,
        "Word.Application",
        17,
        lambda fp, fn: _process_office_textonly(fp, fn, _docx_extract_textonly, "[DOCX]"),
        "[DOCX]",
    )
    if fallback_nodes is not None:
        return fallback_nodes

    return await process_pdf(
        os.path.join(os.path.dirname(file_path), pdf_name),
        images_dir,
        llm_settings,
        shared_llm_url,
        original_filename=pdf_name,
        progress_cb=progress_cb,
        cancel_check=cancel_check,
        keep_vision_alive=keep_vision_alive,
    )
