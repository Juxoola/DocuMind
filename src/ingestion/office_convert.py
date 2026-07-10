"""Конвертация Office-файлов (PPTX, DOCX) в PDF и извлечение текста."""

import asyncio
import logging
import os

import aiofiles
from llama_index.core.schema import TextNode

from src.ingestion.text import process_pdf
from src.ingestion.utils import IngestionCancelled

logger = logging.getLogger(__name__)


# Конвертация Office-файлов в PDF: LibreOffice (приоритет), COM (резерв), текстовый fallback
# ── Поиск LibreOffice: env override → shutil.which → платформо-специфичные пути ──
def _find_soffice():
    import platform
    import shutil

    import config

    # ── env override ──
    env_path = os.getenv("LIBREOFFICE_PATH")
    if env_path and os.path.isfile(env_path):
        return env_path

    # ── shutil.which — работает на всех платформах ──
    found = shutil.which("soffice") or shutil.which("soffice.exe")
    if found:
        return found

    # ── Платформо-специфичные пути ──
    if platform.system() == "Windows":
        local = os.path.join(config.BASE_DIR, "libreoffice", "program", "soffice.exe")
        if os.path.isfile(local):
            return local
        for pf in ["Program Files", "Program Files (x86)"]:
            p = os.path.join("C:\\", pf, "LibreOffice", "program", "soffice.exe")
            if os.path.isfile(p):
                return p
    else:
        for p in ["/usr/bin/soffice", "/usr/lib/libreoffice/program/soffice"]:
            if os.path.isfile(p):
                return p

    return None


async def _convert_via_libreoffice(file_path):
    soffice = _find_soffice()
    if not soffice:
        raise FileNotFoundError("LibreOffice не найден. Установите или скачайте через setup.ps1")
    out_dir = os.path.dirname(file_path)
    cmd = [
        soffice,
        "--headless",
        "--norestore",
        "--convert-to",
        "pdf",
        "--outdir",
        out_dir,
        os.path.abspath(file_path),
    ]
    proc = await asyncio.create_subprocess_exec(
        *cmd,
        stdout=asyncio.subprocess.PIPE,
        stderr=asyncio.subprocess.PIPE,
    )
    _stdout, stderr = await asyncio.wait_for(proc.communicate(), timeout=120)
    pdf_path = os.path.splitext(file_path)[0] + ".pdf"
    if proc.returncode == 0 and await aiofiles.os.path.exists(pdf_path):
        logger.info(f"[DOCX] Сконвертировано в PDF через LibreOffice: {os.path.basename(pdf_path)}")
        try:
            await aiofiles.os.remove(file_path)
        except OSError:
            pass
        return pdf_path
    raise RuntimeError(
        f"LibreOffice конвертация не удалась (code={proc.returncode}): "
        f"{(stderr or b'').decode('utf-8', errors='replace')[:200]}"
    )


# ── Конвертация через COM-интерфейс (PowerPoint/Word) на Windows ──
async def _convert_via_com(file_path, app_name, format_code):

    def _sync_com():
        import pythoncom
        import win32com.client

        pdf_path = os.path.splitext(file_path)[0] + ".pdf"
        app = None
        doc = None
        try:
            pythoncom.CoInitialize()
            app = win32com.client.Dispatch(app_name)
            if app_name == "Powerpoint.Application":
                doc = app.Presentations.Open(os.path.abspath(file_path), WithWindow=False)
            else:
                doc = app.Documents.Open(os.path.abspath(file_path))
            doc.SaveAs(os.path.abspath(pdf_path), format_code)
            if not os.path.exists(pdf_path):
                raise RuntimeError(f"COM {app_name} не создал PDF: {pdf_path}")
            if os.path.exists(file_path):
                os.remove(file_path)
            return pdf_path
        except IngestionCancelled:
            raise
        except Exception as e:
            logger.warning(f"COM-конвертация через {app_name} не удалась: {e}")
            raise
        finally:
            if doc is not None:
                try:
                    doc.Close()
                except Exception:
                    pass
            if app is not None:
                try:
                    app.Quit()
                except Exception:
                    pass
            try:
                pythoncom.CoUninitialize()
            except Exception:
                pass

    return await asyncio.to_thread(_sync_com)


async def _convert_office_to_pdf(file_path, app_name, format_code, textonly_fn, log_prefix):
    file_name = os.path.basename(file_path)
    try:
        pdf_path = await _convert_via_libreoffice(file_path)
    except FileNotFoundError:
        logger.info(f"{log_prefix} LibreOffice не найден, пробую COM-конвертацию...")
        try:
            pdf_path = await _convert_via_com(file_path, app_name, format_code)
        except IngestionCancelled:
            raise
        except Exception as e:
            logger.warning(f"{log_prefix} COM-конвертация тоже не удалась: {e}")
            return None, await textonly_fn(file_path, file_name)
    except IngestionCancelled:
        raise
    except Exception as e:
        logger.warning(f"{log_prefix} LibreOffice конвертация не удалась: {e}")
        try:
            pdf_path = await _convert_via_com(file_path, app_name, format_code)
        except IngestionCancelled:
            raise
        except Exception:
            return None, await textonly_fn(file_path, file_name)
    return os.path.basename(pdf_path), None


def _pptx_extract_textonly(file_path, file_name):
    from pptx import Presentation

    prs = Presentation(file_path)
    nodes = []
    for i, slide in enumerate(prs.slides):
        text = "\n".join([sh.text for sh in slide.shapes if hasattr(sh, "text")])
        if text.strip():
            nodes.append(TextNode(text=text, metadata={"file_name": file_name, "page": i + 1}))
    if nodes:
        logger.info(
            f"[PPTX] Резервный вариант: {len(nodes)} слайдов через python-pptx (без Vision)"
        )
    return nodes


def _docx_extract_textonly(file_path, file_name):
    import docx as _docx

    text = "\n".join([p.text for p in _docx.Document(file_path).paragraphs])
    if text.strip():
        logger.info("[DOCX] Fallback: текст извлечён через python-docx (без Vision)")
        return [TextNode(text=text, metadata={"file_name": file_name})]
    return []


async def _process_office_textonly(file_path, file_name, extract_fn, log_prefix):

    def _sync():
        try:
            return extract_fn(file_path, file_name)
        except Exception as e:
            logger.warning(f"{log_prefix} резервный вариант тоже не удался: {e}")
        return []

    return await asyncio.to_thread(_sync)


# Обработка PPTX/DOCX: конвертация в PDF + Vision или текстовый fallback
# ── Обработка PPTX: конвертация в PDF + Vision или текстовый fallback ──
async def process_pptx(
    file_path,
    images_dir,
    llm_settings=None,
    shared_llm_url=None,
    progress_cb=None,
    cancel_check=None,
    keep_vision_alive=False,
):
    file_name = os.path.basename(file_path)
    pdf_name, fallback_nodes = await _convert_office_to_pdf(
        file_path,
        "Powerpoint.Application",
        32,
        lambda fp, fn: _process_office_textonly(fp, fn, _pptx_extract_textonly, "[PPTX]"),
        "[PPTX]",
    )
    if fallback_nodes is not None:
        return fallback_nodes

    return await process_pdf(
        os.path.join(os.path.dirname(file_path), pdf_name),
        images_dir,
        llm_settings,
        shared_llm_url,
        original_filename=pdf_name,
        progress_cb=progress_cb,
        cancel_check=cancel_check,
        keep_vision_alive=keep_vision_alive,
    )


# ── Обработка DOCX: конвертация в PDF + Vision или текстовый fallback ──
async def process_docx(
    file_path,
    images_dir,
    llm_settings=None,
    shared_llm_url=None,
    progress_cb=None,
    cancel_check=None,
    keep_vision_alive=False,
):
    file_name = os.path.basename(file_path)
    pdf_name, fallback_nodes = await _convert_office_to_pdf(
        file_path,
        "Word.Application",
        17,
        lambda fp, fn: _process_office_textonly(fp, fn, _docx_extract_textonly, "[DOCX]"),
        "[DOCX]",
    )
    if fallback_nodes is not None:
        return fallback_nodes

    return await process_pdf(
        os.path.join(os.path.dirname(file_path), pdf_name),
        images_dir,
        llm_settings,
        shared_llm_url,
        original_filename=pdf_name,
        progress_cb=progress_cb,
        cancel_check=cancel_check,
        keep_vision_alive=keep_vision_alive,
    )
