"""Роутер: загрузка, удаление файлов, метаданные."""

import asyncio
import gc
import logging
import os
import queue
import threading
import time
import uuid as _uuid

import aiofiles
import aiofiles.os
import cv2
import orjson
from fastapi import APIRouter, File, HTTPException, Query, UploadFile
from fastapi.responses import JSONResponse, Response, StreamingResponse

import config

from .shared import (
    _background_tasks,
    _cleanup_ingestion_status,
    _ingestion_lock,
    ingestion_status,
    robust_rmtree,
    safe_filename,
    upload_cancel_flags,
)

logger = logging.getLogger(__name__)
router = APIRouter(tags=["files"])


@router.get("/api/files")
async def get_files(notebook_id: str):
    from routers.notebooks import validate_nb_id

    notebook_id = validate_nb_id(notebook_id)
    paths = config.get_notebook_paths(notebook_id)
    if os.path.exists(paths["data"]):
        files_list = [
            f for f in await aiofiles.os.listdir(paths["data"]) if not f.endswith(".json")
        ]
    else:
        files_list = []
    return {"files": files_list}


@router.post("/api/upload")
async def upload_file(
    file: UploadFile = File(...),
    notebook_id: str = Query(...),
    current_idx: int = Query(1),
    total_count: int = Query(1),
    llm_url: str = Query(None),
    llm_api_key: str | None = None,
    llm_model: str | None = None,
    use_gguf: str | None = None,
    gguf_model_path: str | None = None,
    gguf_mmproj_path: str | None = None,
    vision_model_path: str | None = None,
    vision_mmproj_path: str | None = None,
    vision_temperature: float | None = 0.1,
    vision_ctx_size: int | None = 4096,
    vision_gpu_layers: int | None = -1,
    vision_threads: int | None = 8,
    vision_batch_size: int | None = 512,
    vision_ubatch_size: int | None = 256,
    vision_flash_attn: str | None = "true",
    vision_max_tokens: int | None = 4096,
    vision_repeat_penalty: float | None = 1.2,
    vision_top_p: float | None = 0.9,
    vision_min_p: float | None = 0.05,
    vision_presence_penalty: float | None = 0.0,
    vision_frequency_penalty: float | None = 0.0,
    vision_concurrency: int | None = 1,
    vision_kv_quant: int | None = 2,
    vision_mtp_enabled: bool | None = False,
):
    logger.info(
        f"[API] Новый запрос загрузки для блокнота {notebook_id}. Файл: {file.filename} ({current_idx}/{total_count})"
    )
    _ext = os.path.splitext(file.filename or "")[1].lower()
    if _ext not in config.ALLOWED_UPLOAD_EXTENSIONS:
        raise HTTPException(
            status_code=415,
            detail=f"Тип файла '{_ext}' не поддерживается. Разрешено: {', '.join(sorted(config.ALLOWED_UPLOAD_EXTENSIONS))}.",
        )
    paths = config.get_notebook_paths(notebook_id)
    await aiofiles.os.makedirs(paths["data"], exist_ok=True)
    file_path = os.path.join(paths["data"], safe_filename(file.filename))

    async def save_upload():
        written = 0
        async with aiofiles.open(file_path, "wb") as f:
            while True:
                chunk = await file.read(1024 * 1024)
                if not chunk:
                    break
                written += len(chunk)
                if written > config.UPLOAD_MAX_SIZE_BYTES:
                    try:
                        await aiofiles.os.remove(file_path)
                    except Exception:
                        logger.debug("upload: не удалось удалить недописанный файл")
                    raise HTTPException(
                        status_code=413,
                        detail=f"Файл превысил {config.UPLOAD_MAX_SIZE_MB} МБ во время записи.",
                    )
                await f.write(chunk)

    await save_upload()

    q = queue.Queue()
    effective_llm_url = llm_url
    effective_llm_api_key = llm_api_key
    effective_llm_model = llm_model
    use_gguf_direct = False

    if use_gguf == "true" and gguf_model_path:
        use_gguf_direct = True
        effective_llm_model = os.path.basename(gguf_model_path)

    llm_settings = {
        "llm_url": effective_llm_url,
        "llm_api_key": effective_llm_api_key,
        "llm_model": effective_llm_model,
        "use_gguf_direct": use_gguf_direct,
        "gguf_model_path": gguf_model_path if use_gguf_direct else None,
        "gguf_mmproj_path": gguf_mmproj_path if use_gguf_direct else None,
        "vision_model_path": vision_model_path,
        "vision_mmproj_path": vision_mmproj_path,
        "vision_temperature": vision_temperature,
        "vision_ctx_size": vision_ctx_size,
        "vision_gpu_layers": vision_gpu_layers,
        "vision_threads": vision_threads,
        "vision_batch_size": vision_batch_size,
        "vision_ubatch_size": vision_ubatch_size,
        "vision_flash_attn": vision_flash_attn,
        "vision_max_tokens": vision_max_tokens,
        "vision_repeat_penalty": vision_repeat_penalty,
        "vision_top_p": vision_top_p,
        "vision_min_p": vision_min_p,
        "vision_presence_penalty": vision_presence_penalty,
        "vision_frequency_penalty": vision_frequency_penalty,
        "vision_concurrency": vision_concurrency,
        "vision_kv_quant": vision_kv_quant,
        "vision_mtp_enabled": vision_mtp_enabled,
    }

    task_id = _uuid.uuid4().hex
    q.put({"type": "started", "task_id": task_id, "filename": file.filename})

    def process_task():

        start_time = time.time()
        cancel_event = upload_cancel_flags.setdefault(task_id, threading.Event())
        cancel_event.clear()
        with _ingestion_lock:
            ingestion_status[notebook_id] = {
                "is_uploading": True,
                "progress": 0,
                "batch_progress": (current_idx - 1) / total_count * 100,
                "current_file": current_idx,
                "total_files": total_count,
                "status": "Подготовка...",
                "task_id": task_id,
                "updated_at": time.time(),
            }
        try:
            from src.ingestion import IngestionCancelled
        except ImportError:
            IngestionCancelled = RuntimeError
        try:

            def prog(pct, msg):
                q.put({"type": "progress", "pct": pct, "msg": msg})
                with _ingestion_lock:
                    if notebook_id in ingestion_status:
                        ingestion_status[notebook_id].update({"progress": pct, "status": msg})

            prog(5, "Файл сохранён, подготовка...")
            is_last_in_batch = current_idx >= total_count
            from src.ingestion import ingest_file
            from src.rag.indexing import build_index

            async def _run_pipeline():
                """Единый async-пайплайн: ingest + index на одном event loop.

                Ранее asyncio.run() вызывался дважды (ingest, затем index),
                что ломало httpx.AsyncClient: транспорт привязывался к первому
                loop, а второй asyncio.run() создавал новый — клиент
                оказывался на закрытом loop → 'Event loop is closed'.
                """
                nodes = await ingest_file(
                    file_path,
                    notebook_id,
                    progress_cb=prog,
                    llm_settings=llm_settings,
                    cancel_check=cancel_event.is_set,
                    keep_vision_alive=not is_last_in_batch,
                    keep_whisper_alive=not is_last_in_batch,
                )
                prog(90, "Построение индекса (ChromaDB)...")
                await build_index(nodes, notebook_id)
                return nodes

            nodes = asyncio.run(_run_pipeline())
            from src.rag.retrieval import invalidate_index_cache

            invalidate_index_cache(notebook_id)

            elapsed = time.time() - start_time
            mins = int(elapsed // 60)
            secs = int(elapsed % 60)
            time_str = f"{mins}м {secs}с" if mins > 0 else f"{secs}с"

            if is_last_in_batch:
                logger.info(f"[INGESTION] Пачка завершена. {total_count} файлов обработано.")
                try:
                    from src.gguf.server import unload_all_models

                    asyncio.run(unload_all_models(role="llm"))
                except Exception as llm_err:
                    logger.error(f"[INGESTION] Ошибка выгрузки vision-сервера: {llm_err}")
                try:
                    from src.ingestion import unload_whisper_model

                    asyncio.run(unload_whisper_model())
                except Exception as whisper_err:
                    logger.error(f"[INGESTION] Ошибка выгрузки WhisperX: {whisper_err}")
                try:
                    from src.rag.bm25 import flush_bm25_rebuild

                    asyncio.run(flush_bm25_rebuild(notebook_id))
                except Exception as bm25_err:
                    logger.info(f"[INGESTION] Не удалось форсировать BM25 rebuild: {bm25_err}")
                with _ingestion_lock:
                    ingestion_status[notebook_id] = {
                        "is_uploading": False,
                        "updated_at": time.time(),
                    }
            else:
                with _ingestion_lock:
                    ingestion_status[notebook_id].update(
                        {
                            "batch_progress": current_idx / total_count * 100,
                            "status": f"Готово: {file.filename}",
                        }
                    )

            q.put(
                {
                    "type": "done",
                    "filename": file.filename,
                    "elapsed": time_str,
                    "elapsed_sec": elapsed,
                }
            )
            logger.info(f"[INGESTION] Готово: {file.filename} ({time_str})")
        except IngestionCancelled:
            logger.info(f"[INGESTION] Загрузка отменена пользователем: {file.filename}")
            try:
                from src.gguf.server import kill_stray_servers

                asyncio.run(kill_stray_servers())
            except Exception:
                logger.debug("cancel: не удалось убить llama-server")
            try:
                if os.path.exists(file_path):
                    os.remove(file_path)
            except Exception:
                logger.debug("cancel: не удалось удалить %s", file_path)
            sidecar = os.path.join(os.path.dirname(file_path), f"{file.filename}.json")
            try:
                if os.path.exists(sidecar):
                    os.remove(sidecar)
            except Exception:
                logger.debug("cancel: не удалось удалить sidecar %s", sidecar)
            try:
                images_dir = paths.get("images")
                if images_dir and os.path.exists(images_dir):
                    stem = os.path.splitext(file.filename)[0]
                    for f in os.listdir(images_dir):
                        if f.startswith("p_") or f.startswith("v_") or stem in f:
                            try:
                                os.remove(os.path.join(images_dir, f))
                            except Exception:
                                logger.debug("cancel: не удалось удалить %s", f)
            except Exception:
                logger.debug("cancel: ошибка при чистке images")
            try:
                from src.rag.indexing import get_vector_store

                vector_store = asyncio.run(get_vector_store(notebook_id))
                collection = vector_store._collection
                collection.delete(where={"file_name": file.filename})
            except Exception:
                logger.debug("cancel: не удалось очистить векторные индексы для %s", file.filename)
            with _ingestion_lock:
                ingestion_status[notebook_id] = {
                    "is_uploading": False,
                    "cancelled": True,
                    "updated_at": time.time(),
                }
            q.put({"type": "cancelled", "filename": file.filename})
        except Exception as e:
            logger.error("Ошибка при обработке загрузки", exc_info=True)
            with _ingestion_lock:
                ingestion_status[notebook_id] = {
                    "is_uploading": False,
                    "error": "Внутренняя ошибка обработки файла",
                    "updated_at": time.time(),
                }
            q.put({"type": "error", "msg": "Ошибка обработки файла"})
        finally:
            upload_cancel_flags.pop(task_id, None)
            try:
                from src.ingestion import cleanup_gpu

                cleanup_gpu()
            except Exception:
                logger.debug("finally: не удалось вызвать cleanup_gpu")

    _task = asyncio.create_task(asyncio.to_thread(process_task))
    _background_tasks.add(_task)
    _task.add_done_callback(_background_tasks.discard)

    async def event_generator():
        loop = asyncio.get_running_loop()
        while True:
            msg = await loop.run_in_executor(None, q.get)
            yield f"data: {orjson.dumps(msg).decode()}\n\n"
            if msg["type"] in ("done", "error", "cancelled"):
                break

    return StreamingResponse(event_generator(), media_type="text/event-stream")


@router.post("/api/upload/cancel")
async def cancel_upload(notebook_id: str = Query(...), task_id: str = Query(None)):
    try:
        from src.ingestion import kill_subprocesses

        killed = kill_subprocesses(notebook_id)
        if killed:
            logger.info(f"[CANCEL] Убито {killed} активных subprocess-ов для {notebook_id}")
    except Exception as e:
        logger.error(f"[CANCEL] Ошибка при kill_subprocesses: {e}")
    if task_id:
        evt = upload_cancel_flags.get(task_id)
        if evt is None:
            return {"status": "no_active_upload"}
        evt.set()
        return {"status": "cancel_requested", "task_id": task_id}
    found = False
    for tid, evt in list(upload_cancel_flags.items()):
        evt.set()
        found = True
    return {"status": "cancel_requested" if found else "no_active_upload"}


@router.delete("/api/files/{filename}")
async def delete_file(filename: str, notebook_id: str):
    from routers.notebooks import validate_nb_id

    notebook_id = validate_nb_id(notebook_id)
    filename = safe_filename(filename)
    paths = config.get_notebook_paths(notebook_id)
    file_path = os.path.join(paths["data"], filename)
    if os.path.exists(file_path):
        if filename.lower().endswith((".mp4", ".avi", ".mov")):

            def _release_video_sync(fp):
                cap = cv2.VideoCapture(fp)
                try:
                    cap.get(cv2.CAP_PROP_FPS)
                finally:
                    cap.release()

            await asyncio.to_thread(_release_video_sync, file_path)
        await asyncio.to_thread(gc.collect)

        async def _sync_remove_with_retry(fp: str):
            for i in range(10):
                try:
                    await aiofiles.os.remove(fp)
                    return
                except PermissionError:
                    if i < 4:
                        await asyncio.sleep(1)
                        continue
                    # Fallback: переименовываем и удаляем (обходит лок Windows)
                    try:
                        tmp = fp + f".del{i}"
                        await aiofiles.os.rename(fp, tmp)
                        await aiofiles.os.remove(tmp)
                        return
                    except Exception:
                        if i == 9:
                            raise
                        await asyncio.sleep(1)

        await _sync_remove_with_retry(file_path)
    from src.rag.indexing import get_vector_store

    def _delete_chromadb_entries():
        vs = asyncio.run(get_vector_store(notebook_id))
        vs._collection.delete(where={"file_name": filename})

    await asyncio.to_thread(_delete_chromadb_entries)
    from src.rag.retrieval import invalidate_index_cache

    invalidate_index_cache(notebook_id)
    try:
        from src.bookmarks import mark_stale_for_file

        stale_count = await mark_stale_for_file(notebook_id, filename)
        if stale_count:
            logger.info(
                f"[BOOKMARKS] {stale_count} закладок помечены как stale после удаления {filename}"
            )
    except Exception as e:
        logger.info(f"[BOOKMARKS] Не удалось пометить stale: {e}")
    return {"status": "ok"}


@router.get("/api/source_content")
async def get_source_content(filename: str, notebook_id: str):
    filename = safe_filename(filename)
    ext = os.path.splitext(filename)[1].lower()
    paths = config.get_notebook_paths(notebook_id)
    file_path = os.path.join(paths["data"], filename)

    # Для PDF — чередование текста страниц и описаний изображений
    if ext == ".pdf":
        text = await asyncio.to_thread(
            _build_interleaved_text, file_path, paths["data"], filename
        )
        if text:
            return {"text": text}
    # Для остальных — ChromaDB
    try:
        from src.rag.indexing import get_vector_store
        vector_store = await get_vector_store(notebook_id)
        collection = vector_store._collection
        result = await asyncio.to_thread(collection.get, where={"file_name": filename})
        if result and result.get("documents"):
            full_text = "\n\n---\n\n".join(result["documents"])
            return {"text": full_text}
    except Exception:
        pass
    return {"text": "Содержимое документа не найдено."}


def _extract_pdf_text(file_path: str) -> list[str]:
    """Extract text from PDF page by page."""
    import fitz

    doc = fitz.open(file_path)
    pages = [page.get_text() for page in doc]
    doc.close()
    return pages


def _build_interleaved_text(file_path: str, data_dir: str, filename: str) -> str:
    """Interleave PDF page text with image descriptions from sidecar JSON."""
    pages = _extract_pdf_text(file_path)
    # Загружаем описания из sidecar JSON
    descriptions: dict[int, str] = {}
    json_path = os.path.join(data_dir, f"{filename}.json")
    if os.path.exists(json_path):
        with open(json_path, "rb") as f:
            data = orjson.loads(f.read())
        for frame in data.get("frames", []):
            desc = frame.get("description", "").strip()
            page = frame.get("page")
            if desc and page is not None:
                descriptions[int(page)] = desc
    # Чередуем: текст страницы → описание (если есть)
    # Нумерация страниц в JSON начинается с 1, enumerate — с 0
    parts = []
    for i, page_text in enumerate(pages):
        page_num = i + 1
        if page_text.strip():
            parts.append(page_text)
        if page_num in descriptions:
            parts.append(f"\n\n---\n\n{descriptions[page_num]}\n\n---\n")
    return "\n".join(parts)


def _extract_text_from_file(file_path: str, ext: str) -> str:
    """Extract text directly from file — faster than ChromaDB for large docs."""
    if not os.path.exists(file_path):
        return ""
    try:
        if ext == ".pdf":
            import fitz

            doc = fitz.open(file_path)
            text = "\n".join(page.get_text() for page in doc)
            doc.close()
            return text
        elif ext in (".txt", ".md", ".csv", ".json", ".log"):
            with open(file_path, encoding="utf-8", errors="replace") as f:
                return f.read()
        elif ext in (".docx", ".doc"):
            from docx import Document

            return "\n".join(p.text for p in Document(file_path).paragraphs)
        elif ext in (".pptx", ".ppt"):
            from pptx import Presentation

            texts = []
            for slide in Presentation(file_path).slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        texts.append(shape.text_frame.text)
            return "\n".join(texts)
    except Exception:
        pass
    return ""


_PDF_CSS = """
body { font-family: sans-serif; font-size: 10pt; line-height: 1.5; color: #1a1a1a; }
h1 { font-size: 22pt; margin-bottom: 20px; border-bottom: 2px solid #333; padding-bottom: 8px; }
h2 { font-size: 16pt; margin-top: 24px; color: #222; }
h3 { font-size: 13pt; margin-top: 18px; color: #333; }
p { margin: 6px 0; text-align: justify; }
pre { background: #f4f4f4; padding: 10px; border-radius: 4px; font-size: 9pt; }
code { font-family: monospace; background: #f0f0f0; padding: 1px 4px; border-radius: 2px; }
table { border-collapse: collapse; width: 100%; margin: 12px 0; }
th, td { border: 1px solid #ccc; padding: 6px 8px; text-align: left; font-size: 9pt; }
th { background: #f0f0f0; font-weight: bold; }
blockquote { border-left: 3px solid #999; padding-left: 12px; color: #555; }
"""


def _build_pdf(title: str, text: str) -> bytes:
    """Markdown → HTML → PDF via PyMuPDF Story API."""
    import fitz
    import markdown

    html_body = markdown.markdown(text, extensions=["tables", "fenced_code"])
    html = f"<h1>{title}</h1>\n{html_body}"

    out_path = f"_export_{os.getpid()}.pdf"
    try:
        writer = fitz.DocumentWriter(out_path)
        story = fitz.Story(html=html, user_css=_PDF_CSS)

        def contentfn(positions):
            return html

        def rectfn(page_num, filled):
            return fitz.Rect(0, 0, 595, 842), fitz.Rect(50, 50, 545, 790), fitz.Identity

        story.write_stabilized(writer, contentfn, rectfn, em=10)
        writer.close()

        # write_stabilized не сжимает контент — пересохраняем со сжатием
        doc = fitz.open(out_path)
        compressed_path = f"_export_{os.getpid()}_c.pdf"
        doc.save(compressed_path, garbage=4, deflate=True)
        doc.close()

        with open(compressed_path, "rb") as f:
            pdf_bytes = f.read()


    finally:
        for _f in (out_path, compressed_path):
            try:
                os.unlink(_f)
            except (OSError, UnboundLocalError):
                pass
    return pdf_bytes


@router.get("/api/export_text")
async def export_text(filename: str, notebook_id: str, fmt: str = "txt"):
    """Export extracted text/transcript/description as downloadable file."""
    filename = safe_filename(filename)
    ext = os.path.splitext(filename)[1].lower()
    is_video = ext in (".mp4", ".avi", ".mov", ".mkv")
    is_audio = ext in (".mp3", ".wav", ".m4a", ".aac", ".flac")
    is_image = ext in (".jpg", ".jpeg", ".png", ".webp", ".gif", ".bmp")
    is_media = is_video or is_audio

    text = ""
    stem = os.path.splitext(filename)[0]

    if is_media or is_image:
        # Читаем транскрипт/описание из sidecar JSON
        paths = config.get_notebook_paths(notebook_id)
        json_path = os.path.join(paths["data"], f"{filename}.json")
        if os.path.exists(json_path):
            async with aiofiles.open(json_path, "rb") as f:
                data = orjson.loads(await f.read())
            if is_media:
                transcript = data.get("transcript", [])
                if transcript:
                    lines = []
                    for seg in transcript:
                        start = seg.get("start", 0)
                        m, s = divmod(int(start), 60)
                        h, m = divmod(m, 60)
                        ts = f"{h:02d}:{m:02d}:{s:02d}" if h else f"{m:02d}:{s:02d}"
                        lines.append(f"[{ts}] {seg.get('text', '')}")
                    text = "\n".join(lines)
            elif is_image:
                frames = data.get("frames", [])
                if frames and frames[0].get("description"):
                    text = frames[0]["description"]
        if not text:
            text = f"Описание/транскрипт для {filename} не найдены."
    else:
        # PDF, DOCX, PPT, TXT — читаем из файла
        paths = config.get_notebook_paths(notebook_id)
        file_path = os.path.join(paths["data"], filename)
        if ext == ".pdf":
            text = await asyncio.to_thread(
                _build_interleaved_text, file_path, paths["data"], filename
            )
        else:
            text = await asyncio.to_thread(_extract_text_from_file, file_path, ext)
        if not text:
            text = f"Содержимое {filename} не найдено."

    # Формируем ответ
    base_name = stem
    if fmt == "pdf":
        pdf_bytes = await asyncio.to_thread(_build_pdf, stem, text)
        return Response(
            content=pdf_bytes,
            media_type="application/pdf",
            headers={"Content-Disposition": f'attachment; filename="{base_name}.pdf"'},
        )
    elif fmt == "md":
        content = f"# {stem}\n\n{text}"
        media_type = "text/markdown"
        download_name = f"{base_name}.md"
    else:
        content = text
        media_type = "text/plain; charset=utf-8"
        download_name = f"{base_name}.txt"

    return Response(
        content=content.encode("utf-8"),
        media_type=media_type,
        headers={"Content-Disposition": f'attachment; filename="{download_name}"'},
    )


@router.get("/api/video_metadata")
async def get_video_metadata(filename: str, notebook_id: str):
    filename = safe_filename(filename)
    paths = config.get_notebook_paths(notebook_id)
    json_path = os.path.join(paths["data"], f"{filename}.json")

    async def _async_read_json():
        if not os.path.exists(json_path):
            return None
        async with aiofiles.open(json_path, "rb") as f:
            raw = await f.read()
        return orjson.loads(raw)

    data = await _async_read_json()
    if data is not None:
        return JSONResponse(content=data, headers={"Cache-Control": "public, max-age=300"})
    return {"error": "Метаданные не найдены"}


@router.delete("/api/clear")
async def clear_notebook(notebook_id: str):
    from src.rag.indexing import close_all_clients as _close_all

    async def _clear_data():
        await asyncio.to_thread(_close_all)
        paths = config.get_notebook_paths(notebook_id)
        for d in ("data", "chroma_db", "images"):
            p = paths[d]
            if os.path.exists(p):
                await robust_rmtree(p)
            await asyncio.to_thread(os.makedirs, p, exist_ok=True)

    await _clear_data()
    return {"status": "ok"}


@router.get("/api/ingestion_status")
async def get_ingestion_status(notebook_id: str):
    _cleanup_ingestion_status()
    with _ingestion_lock:
        return ingestion_status.get(notebook_id, {"is_uploading": False})
